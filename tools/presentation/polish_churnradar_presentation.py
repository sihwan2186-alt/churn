from pathlib import Path

from pptx import Presentation


ROOT = Path(__file__).resolve().parents[2]
PRESENTATION_DIR = ROOT / "deliverables" / "presentations"
SOURCE = PRESENTATION_DIR / "ChurnRadar_Detailed_Presentation.pptx"
OUTPUT = PRESENTATION_DIR / "ChurnRadar_Detailed_Presentation_Polished.pptx"
SCRIPT_MD = ROOT / "docs" / "03_presentation" / "PRESENTATION_SCRIPT_POLISHED.md"
SUBMISSION_OUTPUT = ROOT / "submission_package_20260531" / OUTPUT.name
SUBMISSION_SCRIPT = ROOT / "submission_package_20260531" / SCRIPT_MD.name


REPLACEMENTS = {
    "상세 발표본: 데이터, 전처리, 논문 재현, 실험, 운영 해석 전체 포함": "발표용 최종본: 데이터, 전처리, 논문 재현, 실험, 운영 해석 전체 포함",
    "발표 시간이 짧으면 1-18장을 본문, 19장 이후를 백업으로 사용할 수 있다.": "발표 시간이 짧으면 1-18장을 중심으로 진행하고, 19장 이후는 질의응답 자료로 활용합니다.",
    "leakage 없이 어떤 feature를 만들었나?": "정보 누수 없이 어떤 feature를 만들었나?",
    "사용 데이터와 Target": "사용 데이터와 예측 대상",
    "원본 CSV는 B2B 통신사 고객 단위의 정적 CRM snapshot": "시계열 데이터를 찾을 수 없어 정적 CRM 스냅샷 기준으로 진행",
    "정적 CRM snapshot": "정적 CRM 스냅샷",
    "minority class": "소수 클래스",
    "majority class": "다수 클래스",
    "accuracy 단독 사용 불가": "정확도 단독 사용 불가",
    "높은 accuracy보다": "높은 정확도보다",
    "왜 어려운가": "모델링 난점",
    "동일 고객 중복으로 인한 leakage 방지": "동일 고객 중복으로 인한 정보 누수 방지",
    "leakage 위험": "정보 누수 위험",
    "이탈자가 매우 적어 TP 몇 명 차이로 F1이 흔들린다.": "이탈자가 매우 적어 TP 몇 명 차이로 F1이 크게 변동할 수 있다.",
    "월별 행동 변화, 결제 실패, 계약 만료 등 temporal signal이 없다.": "월별 사용량 변화, 결제 이력, 계약 만료 등 시계열 데이터를 찾을 수 없었다.",
    "temporal signal이 없다.": "시계열 데이터를 찾을 수 없었다.",
    "따라서 높은 accuracy보다 재현율과 정밀도의 균형을 함께 봐야 한다.": "따라서 높은 정확도보다 재현율과 정밀도의 균형을 함께 봐야 한다.",
    "recall/precision trade-off를 봐야 한다.": "재현율과 정밀도의 균형을 함께 봐야 한다.",
    "식별자/target/운영 민감 변수는": "식별자/예측 대상/운영 민감 변수는",
    "tree에는 signal, LR에는 noise 가능성": "트리 모델에는 신호, 선형 모델에는 잡음 가능성",
    "CHURN | 이탈 여부 | target | No=0, Yes=1": "CHURN | 이탈 여부 | 예측 대상 | No=0, Yes=1",
    "동일 고객이 train/test에 섞이는 위험 방지": "동일 고객이 학습/테스트에 섞이는 위험 방지",
    "binary classification target": "이진 분류 예측 대상",
    "missing flag + 0 대체": "결측 표시 변수 + 0 대체",
    "이탈 고객이 549명뿐이라 행 삭제는 minority class를 더 줄인다.": "이탈 고객이 549명뿐이라 행 삭제는 소수 클래스를 더 줄인다.",
    "대체값과 결측 flag를 함께 사용해 정보 손실을 줄였다.": "대체값과 결측 표시 변수를 함께 사용해 정보 손실을 줄였다.",
    "shape 확인": "데이터 구조 확인",
    "SVMSMOTE\ntrain only": "SVMSMOTE\n학습 데이터만",
    "중요: resampling과 threshold 선택은 test set을 보지 않고 수행했다.": "재표본화와 임계값 선택은 테스트셋을 보지 않고 수행했습니다.",
    "Leakage 방지 설계": "정보 누수 방지 설계",
    "전처리 성능보다 중요한 것은 test 정보가 train으로 새지 않는 것": "전처리 성능보다 중요한 것은 테스트 정보가 학습 과정에 섞이지 않는 것",
    "test set은 실제 운영 분포처럼 minority 비율 6.5%를 유지했다.": "테스트셋은 실제 운영 분포처럼 이탈 비율 6.5%를 유지했다.",
    "Feature Engineering": "Feature Engineering 설계",
    "모델이 churn 신호를 더 잘 읽는다.": "모델이 이탈 신호를 더 잘 읽는다.",
    "고객 가치와 churn 신호": "고객 가치와 이탈 신호",
    "논문 best는 EasyEnsembleClassifier, F1 0.129.": "논문 최고 성능은 EasyEnsembleClassifier, F1 0.129.",
    "우리 재현 EasyEnsemble F1은 0.128로 거의 동일.": "본 프로젝트의 EasyEnsemble 재현 F1은 0.128로 거의 동일.",
    "LR F1 0.1681은 논문에 없던 추가 모델 조합이므로 직접 우월 주장으로 쓰지 않는다.": "LR F1 0.1681은 논문에 없던 추가 모델 조합이므로 단순 우위 주장으로 사용하지 않습니다.",
    "우리 프로젝트": "본 프로젝트",
    "동일 원천 CSV, 로컬 8,453행": "동일 원천 CSV, 로컬 기준 8,453행",
    "직접 우월 주장": "단순 우위 주장",
    "불균형 데이터에서 baseline, ensemble, categorical model을 함께 확인": "불균형 데이터에서 기준 모델, 앙상블, 범주형 모델을 함께 확인",
    "설명 가능 baseline, 선형 경계 확인": "설명 가능한 기준 모델, 선형 경계 확인",
    "논문 재현, precision 낮음": "논문 재현, 정밀도 낮음",
    "CV 안정성과 recall 균형 좋음": "CV 안정성과 재현율 균형 양호",
    "recall-heavy 후보": "재현율 중심 후보",
    "recall 최고, FP 큼": "재현율 최고, 오탐 많음",
    "결론을 바꿀 정도는 아님": "최종 결론을 바꿀 정도의 개선은 아님",
    "minority class를 실제로 잡는 모델": "소수 클래스를 실제로 탐지하는 모델",
    "하이퍼파라미터와 Threshold 튜닝": "하이퍼파라미터와 임계값 튜닝",
    "모델 점수보다 운영 목적에 맞는 threshold 선택이 중요": "모델 점수보다 운영 목적에 맞는 임계값 선택이 중요",
    "validation F1": "검증 F1",
    "validation threshold 후 test 적용": "검증 임계값 선택 후 테스트 적용",
    "validation recall >= 0.30": "검증 재현율 >= 0.30",
    "우리 추가 모델": "본 프로젝트 추가 모델",
    "paper reference be": "논문 기준",
    "paper reproduction": "논문 재현",
    "main f1 model": "F1 기준",
    "balanced operation": "균형 운영",
    "core recall heavy": "Recall 중심",
    "external paper": "논문 기준",
    "추가 발견": "추가 실험 결과",
    "한 모델이 F1, recall, precision을 모두 이기지는 못했다": "한 모델이 F1, 재현율, 정밀도를 모두 이기지는 못했다",
    "단일 hold-out 결과만으로 우위를 주장하지 않음": "단일 Hold-out 결과만으로 우위를 주장하지 않음",
    "LR은 hold-out F1은 높지만 CV 평균이 0.1309로 낮아져 split 민감성이 있다.": "LR은 Hold-out F1은 높지만 CV 평균이 0.1309로 낮아져 데이터 분할 민감성이 있다.",
    "논문보다 압도적 우위가 아니라 운영 목적별 대안 제시가 타당하다.": "따라서 논문 대비 압도적 성능 향상보다 운영 목적별 대안 제시가 더 타당합니다.",
    "지역 변수는 모델 계열에 따라 signal과 noise가 다르게 작동": "지역 변수는 모델 계열에 따라 신호와 잡음이 다르게 작동",
    "고카디널리티 label encoding은 선형 모델에서 noise가 될 수 있다.": "고카디널리티 label encoding은 선형 모델에서 잡음이 될 수 있다.",
    "어떤 feature 그룹이 어느 모델에서 중요한지 정량화": "어떤 feature 그룹이 어느 모델에서 중요한지 정량적으로 확인",
    "LR: categorical group 제거 시 F1이 0.1681에서 0.0806으로 급락.": "LR은 categorical group 제거 시 F1이 0.1681에서 0.0806으로 급락.",
    "컬럼별 데이터 분리와 단일 컬럼 실험": "컬럼별 데이터 분리와 단일 컬럼 검증",
    "원본 CSV는 수정하지 않고 컬럼별 `feature + CHURN` CSV를 생성했다.": "원본 CSV는 수정하지 않고 컬럼별 feature+CHURN CSV를 생성했다.",
    "단일 컬럼 모델은 성능 상한을 확인하는 빠른 검증용이다.": "단일 컬럼 모델은 변수별 성능 상한을 확인하는 빠른 검증용이다.",
    "각 변수의 역할과 이탈률을 따로 설명하기 위한 보조 데이터": "각 변수의 역할과 이탈률을 따로 확인하기 위한 보조 데이터",
    "컬럼별 성능 한계 확인": "컬럼별 성능 상한 확인",
    "컬럼 하나만 넣으면 F1은 약 0.15 근처에서 한계": "컬럼 하나만 넣으면 F1은 약 0.15 수준에서 한계를 보였다",
    "컬럼별로 좋았던 점과 안 좋았던 점": "컬럼별 강점과 한계",
    "좋았던 점": "강점",
    "안 좋았던 점": "한계",
    "결측률 높아 단독 모델은 약함": "결측률이 높아 단독 모델 성능은 제한적",
    "논문형 core, ZIP top-N, KA 추상화 feature를 실제 성능으로 점검": "논문형 core, ZIP top-N, KA 추상화 feature를 성능으로 점검",
    "최고 F1은 `paper_core_zip_log_ka_abstract + BalancedBagging`의 0.1561.": "최고 F1은 paper_core_zip_log_ka_abstract + BalancedBagging 조합의 0.1561.",
    "기존 최종 LR F1 0.1681을 넘지는 못했다.": "기존 최종 LR F1 0.1681을 상회하지는 못했다.",
    "추가 모델/하이퍼파라미터 실험": "추가 모델과 하이퍼파라미터 실험",
    "BalancedBagging 튜닝은 좋아졌지만 최종 결론은 바꾸지 않음": "BalancedBagging 튜닝 성능은 개선됐지만 최종 결론은 유지",
    "추가 실험 best: `BalancedBagging_tree_depthnone_leaf25`, test F1 0.1605.": "추가 실험 최고 성능은 BalancedBagging_tree_depthnone_leaf25, test F1 0.1605.",
    "기존 최종 모델 `LogisticRegression_SMOTE` F1 0.1681은 유지.": "기존 최종 모델 LogisticRegression_SMOTE의 F1 0.1681은 유지.",
    "전체 평균만 보면 고객군별 실패 패턴을 놓친다": "전체 평균만 보면 고객군별 실패 패턴을 놓칠 수 있다",
    "전용 모델 단독 교체보다 two-stage review가 적합하다.": "전용 모델 단독 교체보다 2단계 검토가 적합하다.",
    "Cost-Sensitive Threshold": "비용 민감 임계값",
    "threshold는 확률값이 아니라 운영 sweep 결과로 해석해야 한다.": "임계값은 확률값 자체가 아니라 운영 조건을 바꿔 본 결과로 해석해야 합니다.",
    "Calibration": "확률 보정",
    "Raw score는 실제 이탈 확률이 아니므로 보정이 필요": "원점수는 실제 이탈 확률이 아니므로 보정이 필요",
    "Feature Importance와 해석 가능성": "Feature Importance와 설명 가능성",
    "모델 합의도 기반 Risk Tier": "모델 합의도 기반 위험 등급",
    "합의도는 영업팀 우선순위 tier로 쓸 수 있다.": "합의도는 영업팀 우선순위 등급으로 쓸 수 있다.",
    "단, vote 0에서도 이탈자가 있어 완전한 rule은 아니다.": "단, vote 0에서도 이탈자가 있어 완전한 규칙은 아니다.",
    "LR의 F1 point estimate는 높지만 신뢰구간이 넓다.": "LR의 F1 점추정치는 높지만 신뢰구간이 넓다.",
    "재학습 기준": "재학습 검토 기준",
    "성능 상한의 핵심은 모델보다 데이터 구조": "성능 개선의 핵심은 모델보다 데이터 구조",
    "이탈 직전 행동 변화 없음": "시계열 데이터 부재로 직전 행동 변화 미반영",
    "모델을 더 추가하는 것보다 temporal feature 확보가 성능 개선의 핵심이다.": "시계열 데이터를 찾을 수 없어 정적 feature로 진행했으며, 향후 시계열 feature 확보가 핵심 개선 방향이다.",
    "따라서 ChurnRadar의 답은 단일 최고 모델이 아니라 운영 목적별 모델 선택 프레임워크다.": "시계열 데이터 부재 조건에서 ChurnRadar의 답은 단일 최고 모델이 아니라 운영 목적별 모델 선택 프레임워크다.",
    "낮은 F1은 실험 실패가 아니라 강한 불균형과 정적 데이터 한계의 결과로 설명한다.": "낮은 F1은 실패가 아니라 강한 불균형과 정적 데이터 구조의 영향으로 해석합니다.",
}


EXACT_REPLACEMENTS = {
    "target": "예측 대상",
}


SCRIPTS = [
    (
        "ChurnRadar",
        "안녕하세요. 이번 발표에서는 B2B 통신사 고객 이탈을 예측하는 ChurnRadar 프로젝트를 설명하겠습니다.\n"
        "핵심은 단순히 가장 높은 점수를 낸 모델을 찾는 것이 아니라, 데이터 제약과 운영 목적에 맞춰 어떤 모델을 선택해야 하는지 정리하는 것입니다.\n"
        "특히 시계열 데이터를 찾을 수 없었기 때문에, 정적 CRM 스냅샷에서 만들 수 있는 feature와 그 한계를 함께 설명하겠습니다.",
    ),
    (
        "상세 발표 흐름",
        "전체 발표 흐름은 예상 질문에 답하는 구조로 구성했습니다.\n"
        "먼저 데이터와 전처리 근거를 설명하고, 그다음 논문 재현 결과와 추가 모델 실험을 비교합니다.\n"
        "시간이 제한되면 18장까지를 본문으로 보고, 이후 슬라이드는 질의응답에서 근거 자료로 사용하겠습니다.",
    ),
    (
        "사용 데이터와 예측 대상",
        "프로젝트 과정에서 월별 사용량, 결제 이력, 계약 만료 같은 시계열 데이터를 찾을 수 없었습니다.\n"
        "그래서 사용 가능한 B2B 통신사 고객 단위의 정적 CRM 스냅샷을 기준으로 이탈 예측을 진행했습니다.\n"
        "전체 8,453개 행 중 이탈 고객은 약 6.5%로 매우 적기 때문에 정확도만 보지 않고 재현율, 정밀도, F1을 함께 확인했습니다.",
    ),
    (
        "원본 컬럼과 사용 여부",
        "이 슬라이드는 어떤 컬럼을 사용했고 어떤 컬럼을 제외했는지 보여줍니다.\n"
        "PID 같은 식별자는 일반화에 도움이 되지 않고 정보 누수 위험이 있어 제외했습니다.\n"
        "반면 매출, 가입자 수, 세그먼트처럼 고객 상태를 설명하는 변수는 모델 입력으로 사용했습니다.",
    ),
    (
        "제외/수정/보존한 데이터",
        "원본 CSV를 직접 수정하지 않고, 모델용 feature만 별도 산출물로 만들었습니다.\n"
        "KA_name은 실명 기반 운영 민감 변수라 기본 모델에서는 제외했고, 연구용으로는 추상화한 형태만 따로 점검했습니다.\n"
        "이 기준은 성능을 조금 높이는 것보다 정보 누수 방지와 실제 운영 가능성을 우선한 결정입니다.",
    ),
    (
        "결측치 처리",
        "결측치는 단순히 삭제하지 않았습니다.\n"
        "특히 이탈 고객 수가 적기 때문에 행 삭제는 소수 클래스를 더 줄여 모델 학습을 어렵게 만듭니다.\n"
        "따라서 대체값과 결측 표시 변수를 함께 사용해, 결측 자체가 가진 운영 신호도 보존했습니다.",
    ),
    (
        "전처리 파이프라인",
        "전처리 과정은 원본 로드부터 보고서 산출까지 재현 가능한 흐름으로 구성했습니다.\n"
        "시계열 데이터가 없었기 때문에 월별 추세 대신 가입자 수, 매출, 세그먼트, 결측 표시 같은 정적 feature를 안정적으로 만들었습니다.\n"
        "중요한 점은 imputation, encoding, scaling, resampling이 모두 학습 데이터 기준으로만 fit되었다는 것입니다.\n"
    ),
    (
        "정보 누수 방지 설계",
        "이 프로젝트에서 가장 조심한 부분은 정보 누수입니다.\n"
        "중앙값 대체, frequency encoding, ZIP top-N grouping처럼 데이터 분포를 학습하는 단계는 모두 train 기준으로만 계산했습니다.\n"
        "이렇게 해야 테스트 성능이 실제 운영 환경에서 기대할 수 있는 성능에 더 가깝게 해석됩니다.",
    ),
    (
        "Feature Engineering 설계",
        "원시 컬럼만 그대로 쓰기보다 이탈 신호를 더 잘 표현하도록 파생 변수를 만들었습니다.\n"
        "가입자 수는 활동성 비율로, 매출은 log와 sqrt 변환으로, 매출과 활동성은 상호작용 변수로 확장했습니다.\n"
        "결론적으로 단일 컬럼보다 파생 feature와 교차 feature가 모델 성능을 설명하는 핵심 근거가 됩니다.",
    ),
    (
        "논문 재현과 비교",
        "먼저 논문의 기준선을 재현해 비교 가능한 출발점을 만들었습니다.\n"
        "논문 최고 성능인 EasyEnsemble F1 0.129와 본 프로젝트 재현값 0.128은 거의 동일합니다.\n"
        "반면 LR 성능은 논문에 없던 추가 실험이므로, 논문보다 무조건 우수하다는 식으로 말하지 않고 별도 기여로 해석합니다.",
    ),
    (
        "논문 방식과 본 프로젝트",
        "이 표는 논문 방식과 본 프로젝트의 차이를 정리한 것입니다.\n"
        "같은 데이터와 유사한 전처리 원칙을 유지하되, 본 프로젝트는 ZIP ablation, 비용 기준 threshold, 운영 관점 비교를 추가했습니다.\n"
        "따라서 재현 비교와 추가 분석을 분리해서 해석하는 것이 중요합니다.",
    ),
    (
        "후보 모델과 사용 이유",
        "불균형 데이터에서는 하나의 모델군만 보기보다 여러 관점의 모델을 비교해야 합니다.\n"
        "Logistic Regression은 설명 가능한 기준 모델이고, EasyEnsemble과 BalancedBagging은 소수 클래스 탐지에 유리한 모델입니다.\n"
        "CatBoost와 XGBoost는 범주형 및 비선형 관계를 확인하기 위해 함께 실험했습니다.",
    ),
    (
        "하이퍼파라미터와 임계값 튜닝",
        "모델 자체의 점수뿐 아니라 운영 목적에 맞는 임계값 선택도 중요합니다.\n"
        "예를 들어 재현율을 높이면 더 많은 이탈 고객을 잡지만 오탐도 함께 증가합니다.\n"
        "그래서 validation 기준으로 threshold를 선택하고, 테스트셋에서는 사후적으로 값을 고르지 않았습니다.",
    ),
    (
        "논문과 비교 가능한 핵심 성능",
        "이 슬라이드가 성능 비교의 중심입니다.\n"
        "EasyEnsemble 재현값은 논문과 거의 일치해 재현성은 확보했습니다.\n"
        "LR은 hold-out F1이 가장 높지만, BalancedBagging과 CatBoost는 운영 목적에 따라 더 높은 재현율을 제공하는 대안입니다.",
    ),
    (
        "Hold-Out 모델 비교",
        "Hold-out 결과를 보면 모델별 장단점이 분명합니다.\n"
        "LR은 F1과 정밀도가 상대적으로 높고, XGBoost와 CatBoost는 재현율이 높지만 오탐이 많습니다.\n"
        "따라서 최종 선택은 단일 점수보다 팀의 접촉 가능 규모와 오탐 비용을 함께 고려해야 합니다.",
    ),
    (
        "5-Fold CV 안정성",
        "단일 Hold-out만 보면 LR이 좋아 보이지만, 5-Fold CV에서는 BalancedBagging과 EasyEnsemble이 더 안정적입니다.\n"
        "이는 데이터 분할에 따라 모델 순위가 달라질 수 있다는 뜻입니다.\n"
        "그래서 이 프로젝트의 결론은 압도적 우위가 아니라 목적별 모델 선택으로 정리하는 것이 타당합니다.",
    ),
    (
        "Billing ZIP Ablation",
        "Billing_ZIP은 모델에 따라 도움이 되기도 하고 방해가 되기도 했습니다.\n"
        "BalancedBagging에서는 ZIP 포함이 F1과 재현율을 높였지만, Logistic Regression에서는 ZIP 제외가 더 좋았습니다.\n"
        "고유값이 많은 범주형 변수는 특히 선형 모델에서 잡음처럼 작동할 수 있음을 보여줍니다.",
    ),
    (
        "Feature Group Ablation",
        "Feature group ablation은 어떤 변수군이 성능에 기여했는지 확인하기 위한 실험입니다.\n"
        "LR에서는 categorical group 제거 시 F1이 크게 떨어졌고, BalancedBagging에서는 interaction group 제거의 영향이 컸습니다.\n"
        "즉 feature 중요도는 모델 구조와 함께 해석해야 합니다.",
    ),
    (
        "컬럼별 데이터 분리와 단일 컬럼 검증",
        "여기서부터는 컬럼별 분석입니다.\n"
        "각 컬럼만 따로 보았을 때 어느 정도 이탈 신호를 갖는지 확인해, 최종 feature 설계가 임의적이지 않다는 근거를 만들었습니다.\n"
        "이 결과는 최종 모델 성능을 대체하는 것이 아니라 feature engineering 방향을 설명하는 보조 근거입니다.",
    ),
    (
        "컬럼별 CSV 생성 결과",
        "컬럼별 CSV와 요약 파일은 발표에서 변수별 의미를 설명하기 위한 보조 자료입니다.\n"
        "예를 들어 단일 변수 설명, 범주값별 이탈률, 단일 컬럼 모델 성능을 각각 따로 확인할 수 있게 했습니다.\n"
        "이렇게 분리해 두면 질의응답에서 특정 컬럼의 역할을 더 명확하게 설명할 수 있습니다.",
    ),
    (
        "단일 컬럼 모델 스크리닝",
        "단일 컬럼만 사용하면 최고 F1이 약 0.15 수준에서 멈춥니다.\n"
        "가입자 수와 매출 관련 컬럼이 상대적으로 강하지만, 단독으로는 충분한 예측력을 제공하지 못합니다.\n"
        "이 결과는 여러 feature를 결합하고 변환해야 하는 이유를 뒷받침합니다.",
    ),
    (
        "컬럼별 강점과 한계",
        "이 표는 단일 컬럼 실험을 최종 feature 설계와 연결한 것입니다.\n"
        "가입자 수와 매출은 중요한 신호지만, 규모 자체만으로는 정밀도가 낮습니다.\n"
        "그래서 비율, 변환, 상호작용 feature로 확장해 모델이 더 안정적으로 신호를 읽도록 했습니다.",
    ),
    (
        "Paper/KA Ablation Variants",
        "이 실험은 논문형 core feature와 KA 추상화 feature를 실제 성능으로 점검한 것입니다.\n"
        "최고 F1은 0.1561로 의미 있는 개선이 있었지만 기존 LR F1 0.1681을 넘지는 못했습니다.\n"
        "따라서 KA 추상화 feature는 최종 결론을 바꾸는 근거라기보다 보조 분석으로 활용하는 것이 적절합니다.",
    ),
    (
        "추가 모델과 하이퍼파라미터 실험",
        "추가 튜닝에서는 BalancedBagging의 성능이 일부 개선되었습니다.\n"
        "다만 최고 F1이 0.1605로 기존 LR 0.1681을 넘지는 않았습니다.\n"
        "그래서 최종 모델 결론은 유지하되, BalancedBagging은 재현율 중심 운영 후보로 설명할 수 있습니다.",
    ),
    (
        "CRM Segment 분석",
        "전체 평균만 보면 고객군별로 어떤 문제가 생기는지 놓칠 수 있습니다.\n"
        "mid/high value에서는 재현율은 높지만 정밀도가 낮아 오탐 관리가 중요하고, low value에서는 놓치는 이탈자가 상대적으로 많습니다.\n"
        "이 결과는 segment별 운영 전략이 필요하다는 점을 보여줍니다.",
    ),
    (
        "High-Value 전용 모델 실험",
        "고가치 고객만 따로 학습한 모델은 정밀도와 PR-AUC가 개선됐지만 재현율이 크게 낮아졌습니다.\n"
        "즉 고가치 전용 모델 하나로 전체 모델을 대체하기에는 위험합니다.\n"
        "더 적합한 방식은 전체 모델로 후보를 넓게 잡고, 고가치 고객은 2단계 검토로 정밀하게 보는 것입니다.",
    ),
    (
        "비용 민감 임계값",
        "비용 구조를 넣으면 모델 선택 기준이 달라집니다.\n"
        "이탈 방지 가치가 큰 시나리오에서는 재현율을 극대화하는 쪽이 유리하지만, 캠페인 비용이 커지면 기대가치가 음수가 될 수 있습니다.\n"
        "따라서 threshold는 확률 자체가 아니라 운영 조건별 의사결정 결과로 해석해야 합니다.",
    ),
    (
        "비즈니스 임팩트 시나리오",
        "이 슬라이드는 모델 성능을 비즈니스 의사결정으로 연결한 요약입니다.\n"
        "같은 모델이라도 접촉 가능 고객 수, 캠페인 비용, 이탈 방지 가치에 따라 실제 효과가 달라집니다.\n"
        "따라서 운영팀이 쓸 수 있는 기준은 점수 하나가 아니라 비용과 예산을 반영한 선택 기준입니다.",
    ),
    (
        "Top-k 예산 전략",
        "실제 캠페인은 threshold보다 상위 몇 퍼센트를 접촉할지로 운영하는 경우가 많습니다.\n"
        "예산이 작으면 LR이나 EasyEnsemble이 효율적이고, 예산이 커질수록 BalancedBagging 계열의 재현율 장점이 커집니다.\n"
        "이 슬라이드는 모델 선택을 캠페인 예산과 연결해 설명하는 근거입니다.",
    ),
    (
        "확률 보정",
        "모델의 raw score는 실제 이탈 확률이 아닙니다.\n"
        "예를 들어 score 0.4라고 해서 이탈 확률이 40%라는 뜻은 아니므로, 확률처럼 사용하려면 보정이 필요합니다.\n"
        "Platt 보정 후 Brier와 ECE가 낮아진 것은 운영에서 점수를 더 안정적으로 해석할 수 있음을 보여줍니다.",
    ),
    (
        "Feature Importance와 설명 가능성",
        "해석 방법마다 상위 변수 순위는 다를 수 있습니다.\n"
        "하지만 논문 SHAP, permutation importance, LR coefficient 모두 활동성이나 매출 관련 신호가 중요하다는 방향으로 수렴합니다.\n"
        "따라서 모델 결과는 단순 블랙박스가 아니라 고객 활동성과 수익성 관점에서 설명할 수 있습니다.",
    ),
    (
        "모델 합의도 기반 위험 등급",
        "여러 모델이 동시에 이탈 위험을 경고한 고객군은 실제 이탈률도 높았습니다.\n"
        "특히 8개 모델 모두가 경고한 그룹은 전체 평균보다 약 1.9배 높은 이탈률을 보였습니다.\n"
        "다만 경고가 없는 그룹에서도 이탈자는 존재하므로, 합의도는 절대 규칙이 아니라 우선순위 등급으로 사용해야 합니다.",
    ),
    (
        "Bootstrap CI와 McNemar 검정",
        "성능 차이를 단일 점수만으로 판단하지 않기 위해 신뢰구간과 McNemar 검정을 확인했습니다.\n"
        "LR은 F1 점추정치가 높지만 신뢰구간이 넓고, recall 중심 모델들은 이탈 탐지 측면의 일관성이 있습니다.\n"
        "즉 모델 간 차이는 점수 차이뿐 아니라 어떤 고객을 맞히고 틀리는지의 차이로 봐야 합니다.",
    ),
    (
        "재현 자동화와 MLOps",
        "제출 이후에도 같은 결과를 재현할 수 있도록 자동화 구조를 만들었습니다.\n"
        "pytest로 데이터 스키마를 확인하고, preprocess와 phase 3-8 실험을 재실행하며, n8n workflow로 전체 흐름을 묶었습니다.\n"
        "운영 단계에서는 drift check와 재학습 검토 기준을 추가해 프로젝트를 지속적으로 관리할 수 있습니다.",
    ),
    (
        "한계와 향후 개선",
        "가장 큰 한계는 시계열 데이터를 찾을 수 없어 정적 CRM 스냅샷 기준으로 진행했다는 점입니다.\n"
        "그래서 이탈 직전의 사용량 변화, 결제 실패, VOC, 계약 만료 정보가 없고, 모델 성능에는 구조적인 상한이 있습니다.\n"
        "향후에는 시계열 feature와 외부 검증 데이터를 확보하는 것이 모델을 추가하는 것보다 더 중요한 개선 방향입니다.",
    ),
    (
        "최종 결론",
        "최종적으로 논문 EasyEnsemble baseline은 F1 0.128로 재현했습니다.\n"
        "Hold-out에서는 LR이 가장 높았지만, CV 안정성과 운영 재현율 관점에서는 BalancedBagging과 EasyEnsemble도 중요한 대안입니다.\n"
        "시계열 데이터를 찾을 수 없었던 조건에서는 점수 하나보다 데이터 한계를 인정하고 운영 목적별 모델 선택 기준을 제시하는 것이 핵심입니다.",
    ),
]


def iter_text_frames(shape):
    if getattr(shape, "has_text_frame", False):
        yield shape.text_frame
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                yield cell.text_frame


def replace_in_text_frame(text_frame):
    count = 0
    for old, new in sorted(REPLACEMENTS.items(), key=lambda item: len(item[0]), reverse=True):
        for paragraph in text_frame.paragraphs:
            replaced_in_runs = False
            for run in paragraph.runs:
                if old in run.text:
                    run.text = run.text.replace(old, new)
                    count += 1
                    replaced_in_runs = True
            if not replaced_in_runs and old in paragraph.text:
                paragraph.text = paragraph.text.replace(old, new)
                count += 1
    for paragraph in text_frame.paragraphs:
        exact = paragraph.text.strip()
        if exact in EXACT_REPLACEMENTS:
            if paragraph.runs:
                for run in paragraph.runs:
                    run.text = run.text.replace(exact, EXACT_REPLACEMENTS[exact])
            else:
                paragraph.text = EXACT_REPLACEMENTS[exact]
            count += 1
    return count


def first_visible_title(slide):
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            text = shape.text.strip()
            if text and not text.startswith("ChurnRadar |"):
                return text.splitlines()[0]
    return "Untitled"


def set_notes(slide, script):
    notes = slide.notes_slide
    frame = notes.notes_text_frame
    frame.clear()
    frame.text = script


def write_script_markdown(prs):
    lines = [
        "# ChurnRadar 발표 대본",
        "",
        "각 슬라이드의 발표자 노트에도 같은 대본을 삽입했습니다.",
        "",
    ]
    for index, (title, script) in enumerate(SCRIPTS, start=1):
        slide_title = first_visible_title(prs.slides[index - 1])
        lines.append(f"## Slide {index}. {slide_title}")
        lines.append("")
        lines.append(script)
        lines.append("")
    SCRIPT_MD.parent.mkdir(parents=True, exist_ok=True)
    SCRIPT_MD.write_text("\n".join(lines), encoding="utf-8")
    if SUBMISSION_SCRIPT.parent.exists():
        SUBMISSION_SCRIPT.write_text("\n".join(lines), encoding="utf-8")


def main():
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation(SOURCE)
    if len(prs.slides) != len(SCRIPTS):
        raise RuntimeError(f"Expected {len(SCRIPTS)} scripts, found {len(prs.slides)} slides")

    replacement_count = 0
    for slide_index, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            for text_frame in iter_text_frames(shape):
                replacement_count += replace_in_text_frame(text_frame)
        set_notes(slide, SCRIPTS[slide_index][1])

    prs.save(OUTPUT)
    if SUBMISSION_OUTPUT.parent.exists():
        prs.save(SUBMISSION_OUTPUT)
    write_script_markdown(prs)

    print(f"Saved PPTX: {OUTPUT}")
    if SUBMISSION_OUTPUT.exists():
        print(f"Saved submission PPTX: {SUBMISSION_OUTPUT}")
    print(f"Saved script: {SCRIPT_MD}")
    if SUBMISSION_SCRIPT.exists():
        print(f"Saved submission script: {SUBMISSION_SCRIPT}")
    print(f"Slides: {len(prs.slides)}")
    print(f"Text replacements applied: {replacement_count}")


if __name__ == "__main__":
    main()
