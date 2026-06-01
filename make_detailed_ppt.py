from pathlib import Path
from typing import Iterable

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.util import Inches, Pt

from make_final_ppt import (
    ACCENT,
    ACCENT_2,
    ACCENT_3,
    ASSETS,
    FONT,
    LIGHT,
    MUTED,
    PHASE5B,
    PHASE6,
    PHASE8,
    ROOT,
    TEXT,
    TITLE,
    WHITE,
    add_band,
    add_bullets,
    add_footer,
    add_picture_fit,
    add_text,
    blank_slide,
    set_run,
)


OUTPUT = ROOT / "ChurnRadar_Detailed_Presentation.pptx"
SUBMISSION_OUTPUT = ROOT / "submission_package_20260531" / "ChurnRadar_Detailed_Presentation.pptx"

PHASE3B = ROOT / "processed" / "phase_3b_differentiation"
PHASE4 = ROOT / "processed" / "phase_4_paper_comparison"
PHASE7 = ROOT / "processed" / "phase_7_next_experiments"
COLUMN_SPLIT = ROOT / "processed" / "column_split_datasets"
ADDITIONAL = ROOT / "processed" / "additional_experiments"


def fmt(value, digits: int = 3) -> str:
    if pd.isna(value):
        return ""
    if isinstance(value, str):
        return value
    if abs(float(value)) >= 1000 and float(value).is_integer():
        return f"{int(value):,}"
    return f"{float(value):.{digits}f}"


def pct(value, digits: int = 1) -> str:
    if pd.isna(value):
        return ""
    return f"{float(value) * 100:.{digits}f}%"


def short_model(name: str) -> str:
    mapping = {
        "LogisticRegression_SMOTE": "LR SMOTE",
        "BalancedBagging_original": "BalancedBagging",
        "CatBoost_native_categorical": "CatBoost native",
        "XGBoost_SMOTE": "XGBoost",
        "EasyEnsembleClassifier": "Paper EasyEnsemble",
        "EasyEnsemble_original": "EasyEnsemble",
    }
    return mapping.get(str(name), str(name).replace("_", " "))


def short_case(name: str) -> str:
    mapping = {
        "LR_no_zip_f1": "LR no ZIP",
        "BalancedBagging_with_zip": "BalancedBagging",
        "CatBoost_native_with_zip": "CatBoost native",
        "XGBoost_with_zip": "XGBoost",
        "EasyEnsemble_with_zip": "EasyEnsemble",
    }
    return mapping.get(str(name), str(name).replace("_", " "))


def short_purpose(name: str) -> str:
    mapping = {
        "main_f1_model": "F1 기준",
        "balanced_operation_model": "균형 운영",
        "core_recall_heavy_model": "Recall 중심",
        "extended_recall_heavy_model": "Recall 확장",
        "cost_optimized_operating_point": "비용 최적",
        "paper_reference_best": "논문 기준",
        "paper_reproduction_model": "재현 모델",
    }
    return mapping.get(str(name), str(name).replace("_", " "))


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    add_text(slide, title, 0.55, 0.22, 12.2, 0.42, size=22, bold=True, color=TITLE)
    if subtitle:
        add_text(slide, subtitle, 0.58, 0.69, 11.85, 0.32, size=10, color=MUTED)


def add_section_slide(prs: Presentation, page: int, title: str, subtitle: str, bullets: Iterable[str]) -> None:
    slide = blank_slide(prs)
    add_band(slide, 0, 0, 13.333, 7.5, RGBColor(235, 242, 248))
    add_text(slide, title, 0.75, 1.0, 11.7, 0.75, size=36, bold=True, color=TITLE)
    add_text(slide, subtitle, 0.82, 1.9, 10.8, 0.5, size=17, color=TEXT)
    add_bullets(slide, bullets, 1.0, 3.1, 11.0, 2.2, size=17)
    add_footer(slide, page)


def add_small_table(
    slide,
    rows: list[list[str]],
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: int = 8,
    header_fill: RGBColor = TITLE,
    first_col_width: float | None = None,
) -> None:
    row_count = len(rows)
    col_count = len(rows[0])
    table_shape = slide.shapes.add_table(
        row_count,
        col_count,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    table = table_shape.table
    if first_col_width:
        table.columns[0].width = Inches(first_col_width)
        rest = max(width - first_col_width, 0.1) / max(col_count - 1, 1)
        for idx in range(1, col_count):
            table.columns[idx].width = Inches(rest)
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(value)
            cell.margin_left = Inches(0.04)
            cell.margin_right = Inches(0.04)
            cell.margin_top = Inches(0.025)
            cell.margin_bottom = Inches(0.025)
            cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if c else PP_ALIGN.LEFT
                for run in paragraph.runs:
                    set_run(
                        run,
                        size=font_size,
                        bold=(r == 0),
                        color=WHITE if r == 0 else TEXT,
                    )
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(248, 250, 252)


def add_process_lane(slide, steps: list[tuple[str, str]], left: float, top: float, width: float) -> None:
    box_w = width / len(steps) - 0.08
    for idx, (label, desc) in enumerate(steps):
        x = left + idx * (box_w + 0.08)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(top), Inches(box_w), Inches(1.25))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(238, 246, 250) if idx % 2 == 0 else RGBColor(247, 250, 247)
        shape.line.color.rgb = RGBColor(214, 224, 235)
        frame = shape.text_frame
        frame.clear()
        frame.margin_left = Inches(0.08)
        frame.margin_right = Inches(0.08)
        frame.margin_top = Inches(0.06)
        frame.margin_bottom = Inches(0.06)
        p = frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        set_run(run, size=10, bold=True, color=TITLE)
        p2 = frame.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = desc
        set_run(run2, size=8, color=TEXT)
        if idx < len(steps) - 1:
            add_text(slide, ">", x + box_w + 0.01, top + 0.45, 0.08, 0.2, size=13, bold=True, color=MUTED, align=PP_ALIGN.CENTER)


def rows_from_df(df: pd.DataFrame, columns: list[tuple[str, str]], n: int, digits: int = 3) -> list[list[str]]:
    rows = [[header for header, _ in columns]]
    for _, row in df.head(n).iterrows():
        out = []
        for _, col in columns:
            value = row[col]
            if col.lower().endswith("rate") or col in {"recall", "precision", "f1", "pr_auc", "test_f1", "test_recall", "test_precision"}:
                out.append(fmt(value, digits))
            else:
                out.append(fmt(value, digits))
        rows.append(out)
    return rows


def build_detailed_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    page = 1

    final = pd.read_csv(ROOT / "final_model_summary.csv")
    single = pd.read_csv(COLUMN_SPLIT / "05_single_column_best_models.csv")
    cv = pd.read_csv(PHASE4 / "phase_4_cv_summary.csv")
    zip_ablation = pd.read_csv(PHASE3B / "experiment_a_billing_zip_summary.csv")
    feature_ablation = pd.read_csv(PHASE3B / "experiment_b_feature_group_ablation.csv")
    segment = pd.read_csv(PHASE3B / "experiment_c_segment_bucket_summary.csv")
    high_value = pd.read_csv(PHASE3B / "experiment_c_high_value_submodel.csv")
    cost_best = pd.read_csv(PHASE3B / "experiment_d_cost_threshold_best.csv")
    additional = pd.read_csv(ADDITIONAL / "additional_top25_summary.csv")
    paper_ablation = pd.read_csv(PHASE7 / "paper_ablation_top30.csv")
    topk = pd.read_csv(PHASE6 / "phase6_topk_budget_curve.csv")
    calibration = pd.read_csv(PHASE6 / "phase6_calibration_metrics.csv")
    metric_ci = pd.read_csv(PHASE8 / "bootstrap_metric_ci.csv")

    slide = blank_slide(prs)
    add_band(slide, 0, 0, 13.333, 7.5, RGBColor(235, 242, 248))
    add_text(slide, "ChurnRadar", 0.65, 0.72, 12.0, 0.8, size=42, bold=True, color=TITLE)
    add_text(slide, "상세 발표본: 데이터, 전처리, 논문 재현, 실험, 운영 해석 전체 포함", 0.72, 1.58, 11.6, 0.4, size=18, color=TEXT)
    add_bullets(
        slide,
        [
            "B2B 통신사 고객 이탈 예측 데이터의 컬럼별 역할과 제외/수정 사유 설명",
            "논문 EasyEnsemble baseline 재현 후 모델/feature/threshold 실험으로 확장",
            "최종 결론: 단일 최고 모델이 아니라 운영 목적별 모델 선택 프레임워크",
        ],
        0.8,
        2.55,
        6.3,
        2.0,
        size=16,
    )
    add_picture_fit(slide, ASSETS / "01_model_metric_comparison.png", 7.1, 2.25, 5.5, 3.3)
    add_footer(slide, page)
    page += 1

    slide = blank_slide(prs)
    add_title(slide, "상세 발표 흐름", "질문을 예상하고 답하는 순서로 구성")
    add_small_table(
        slide,
        [
            ["구간", "슬라이드", "핵심 질문"],
            ["데이터 이해", "3-6", "무슨 데이터이고 무엇을 빼거나 수정했나?"],
            ["전처리", "7-9", "leakage 없이 어떤 feature를 만들었나?"],
            ["논문 비교", "10-14", "논문과 비교 가능한 값은 무엇인가?"],
            ["모델/튜닝", "15-23", "어떤 실험과 하이퍼파라미터 튜닝을 했나?"],
            ["컬럼별 분석", "19-21", "각 컬럼만 넣으면 무엇이 좋고 안 좋은가?"],
            ["운영 해석", "24-31", "비용, segment, top-k, calibration에서 무엇을 배웠나?"],
            ["검증/결론", "32-35", "통계적으로 얼마나 믿을 수 있고 한계는 무엇인가?"],
        ],
        0.8,
        1.25,
        11.8,
        3.0,
        font_size=10,
        first_col_width=2.0,
    )
    add_bullets(
        slide,
        [
            "17장 요약본보다 방어 자료 성격이 강한 상세본이다.",
            "발표 시간이 짧으면 1-18장을 본문, 19장 이후를 백업으로 사용할 수 있다.",
        ],
        1.05,
        5.0,
        11.0,
        1.0,
        size=14,
    )
    add_footer(slide, page)
    page += 1

    slide = blank_slide(prs)
    add_title(slide, "사용 데이터와 Target", "원본 CSV는 B2B 통신사 고객 단위의 정적 CRM snapshot")
    add_small_table(
        slide,
        [
            ["항목", "값", "발표 해석"],
            ["원본 파일", "Baza customer Telecom v2.csv", "프로젝트 전체의 입력 데이터"],
            ["원본 크기", "8,453 x 14", "14개 컬럼에는 target CHURN 포함"],
            ["PID 중복 제거 후", "8,436 rows", "동일 고객 중복으로 인한 leakage 방지"],
            ["CHURN=Yes", "549", "minority class"],
            ["CHURN=No", "7,904", "majority class"],
            ["이탈 비율", "약 6.5%", "accuracy 단독 사용 불가"],
        ],
        0.75,
        1.15,
        6.35,
        2.8,
        font_size=9,
        first_col_width=1.55,
    )
    add_band(slide, 7.45, 1.25, 4.9, 3.0, RGBColor(238, 246, 250))
    add_text(slide, "왜 어려운가", 7.75, 1.55, 4.2, 0.35, size=19, bold=True, color=TITLE)
    add_bullets(
        slide,
        [
            "이탈자가 매우 적어 TP 몇 명 차이로 F1이 흔들린다.",
            "월별 행동 변화, 결제 실패, 계약 만료 등 temporal signal이 없다.",
            "따라서 높은 accuracy보다 recall/precision trade-off를 봐야 한다.",
        ],
        7.75,
        2.05,
        4.25,
        1.6,
        size=12,
    )
    add_footer(slide, page)
    page += 1

    slide = blank_slide(prs)
    add_title(slide, "원본 컬럼과 사용 여부", "식별자/target/운영 민감 변수는 분리하고, 고객 행동/매출 변수 중심으로 학습")
    add_small_table(
        slide,
        [
            ["컬럼", "의미", "사용", "이유"],
            ["PID", "고객 식별자", "제외", "일반화 불가, leakage 위험"],
            ["KA_name", "담당자명", "기본 제외", "개인/조직 변화 민감, 연구용 추상화만 별도"],
            ["CRM_PID_Value_Segment", "고객 가치 등급", "사용", "segment별 위험 차이 반영"],
            ["EffectiveSegment", "사업/규모 세그먼트", "사용", "SOHO/VSE/SME 등 고객군 차이"],
            ["Billing_ZIP", "청구 우편번호", "비교", "tree에는 signal, LR에는 noise 가능성"],
            ["가입자 수", "활성/비활성/정지/전체", "사용", "서비스 이용 상태"],
            ["매출 변수", "모바일/유선/총매출/ARPU", "사용", "고객 가치와 churn 신호"],
            ["CHURN", "이탈 여부", "target", "No=0, Yes=1"],
        ],
        0.55,
        1.12,
        12.25,
        4.15,
        font_size=8,
        first_col_width=2.05,
    )
    add_text(slide, "핵심: 무엇을 뺐는지보다 왜 뺐는지를 설명해야 방어가 된다.", 0.75, 6.0, 12.0, 0.35, size=14, bold=True, color=ACCENT)
    add_footer(slide, page)
    page += 1

    slide = blank_slide(prs)
    add_title(slide, "제외/수정/보존한 데이터", "원본은 보존하고 모델용 feature만 목적에 맞게 변환")
    add_small_table(
        slide,
        [
            ["항목", "처리", "이유"],
            ["PID", "중복 제거 키로만 사용 후 제거", "동일 고객이 train/test에 섞이는 위험 방지"],
            ["KA_name", "기본 모델 제외", "담당자명은 운영 민감 변수이며 일반화 취약"],
            ["KA_name 연구용", "code type / premium interaction 추상화", "실명 제거 후 signal 가능성만 점검"],
            ["Sliver", "Silver로 통합", "명백한 오타 정정"],
            ["Billing_ZIP", "포함/제외/top-50 variant", "고카디널리티 지역 변수의 효과 검증"],
            ["ARPU 결측", "TotalRevenue/Total_SUBs 보정 후 median", "수익성 feature 보존"],
            ["CHURN", "No=0, Yes=1", "binary classification target"],
        ],
        0.75,
        1.18,
        11.85,
        3.15,
        font_size=9,
        first_col_width=2.2,
    )
    add_bullets(
        slide,
        [
            "모델 성능을 높이기 위한 임의 삭제보다 leakage와 운영 가능성을 우선했다.",
            "원본 CSV는 수정하지 않고, 전처리 산출물만 별도 폴더에 생성했다.",
        ],
        1.0,
        5.05,
        11.3,
        1.0,
        size=14,
    )
    add_footer(slide, page)
    page += 1

    slide = blank_slide(prs)
    add_title(slide, "결측치 처리", "결측 자체가 이탈 신호일 수 있어 flag로 보존")
    add_small_table(
        slide,
        [
            ["컬럼", "결측률", "처리", "사용한 이유"],
            ["Suspended_subscribers", "약 95.84%", "missing flag + 0 대체", "정지 가입자 없음/기록 부재 가능성"],
            ["Not_Active_subscribers", "약 49.08%", "missing flag + 0 대체", "비활성 상태 자체가 위험 신호 가능"],
            ["CRM_PID_Value_Segment", "약 0.06%", "Unknown", "행 삭제 시 minority class 손실"],
            ["Billing_ZIP", "약 0.02%", "train median 또는 제외", "지역 신호 검증"],
            ["ARPU", "약 0.01%", "보정 후 train median", "수익성 feature 보존"],
        ],
        0.7,
        1.25,
        12.0,
        2.35,
        font_size=9,
        first_col_width=2.35,
    )
    add_text(slide, "삭제하지 않은 이유", 0.9, 4.3, 3.0, 0.3, size=18, bold=True, color=TITLE)
    add_bullets(
        slide,
        [
            "이탈 고객이 549명뿐이라 행 삭제는 minority class를 더 줄인다.",
            "결측이 무작위가 아니라 운영 상태를 반영할 수 있다.",
            "대체값과 결측 flag를 함께 사용해 정보 손실을 줄였다.",
        ],
        1.0,
        4.75,
        11.5,
        1.25,
        size=14,
    )
    add_footer(slide, page)
    page += 1

    slide = blank_slide(prs)
    add_title(slide, "전처리 파이프라인", "원본 CSV에서 모델 입력과 평가표까지 재현 가능한 흐름")
    add_process_lane(
        slide,
        [
            ("1. Load", "원본 CSV\nshape 확인"),
            ("2. Dedup", "PID 기준\n17건 제거"),
            ("3. Split", "80:20\nstratified"),
            ("4. Impute", "train 기준\n대체값 fit"),
            ("5. Encode", "label/frequency\ntrain 기준"),
            ("6. Scale", "StandardScaler\ntrain 기준"),
        ],
        0.65,
        1.35,
        12.0,
    )
    add_process_lane(
        slide,
        [
            ("7. Feature", "rate/log/sqrt\ninteraction"),
            ("8. Resample", "SVMSMOTE\ntrain only"),
            ("9. Model", "LR/ensemble\nCatBoost/XGB"),
            ("10. Tune", "validation\nthreshold"),
            ("11. Test", "test 1회\n최종 평가"),
            ("12. Report", "CSV/PNG/PPT\n산출"),
        ],
        0.65,
        3.4,
        12.0,
    )
    add_text(slide, "중요: resampling과 threshold 선택은 test set을 보지 않고 수행했다.", 0.85, 5.9, 11.5, 0.35, size=14, bold=True, color=ACCENT)
    add_footer(slide, page)
    page += 1

    slide = blank_slide(prs)
    add_title(slide, "Leakage 방지 설계", "전처리 성능보다 중요한 것은 test 정보가 train으로 새지 않는 것")
    add_small_table(
        slide,
        [
            ["작업", "fit 위치", "test 처리"],
            ["Imputation median", "train only", "train median으로 transform"],
            ["Label/frequency encoding", "train only", "unseen은 Unknown/0 frequency"],
            ["ZIP top-N grouping", "train frequency", "train top-N 기준 적용"],
            ["StandardScaler", "train only", "transform만 적용"],
            ["SVMSMOTE", "train only", "test에는 적용하지 않음"],
            ["Threshold 선택", "train 내부 validation", "test에는 1회 적용"],
            ["KA target encoding", "train fold/leave-one-out", "test는 train mapping만 적용"],
        ],
        0.75,
        1.15,
        11.85,
        3.2,
        font_size=9,
        first_col_width=2.7,
    )
    add_bullets(
        slide,
        [
            "논문과 비교 가능한 재현성을 위해 random_state와 split 원칙을 고정했다.",
            "test set은 실제 운영 분포처럼 minority 비율 6.5%를 유지했다.",
        ],
        1.0,
        5.05,
        11.1,
        1.0,
        size=14,
    )
    add_footer(slide, page)
    page += 1

    slide = blank_slide(prs)
    add_title(slide, "Feature Engineering", "원시 매출/가입자 수를 활동성, 수익성, 상호작용 신호로 확장")
    add_small_table(
        slide,
        [
            ["그룹", "예시", "왜 사용했나"],
            ["원시 매출", "AvgMobileRevenue, TotalRevenue, ARPU", "고객 가치와 이탈 비용 반영"],
            ["가입자 수", "Active, Not Active, Suspended, Total_SUBs", "서비스 사용 상태 반영"],
            ["비율", "active_rate, inactive_rate, dormant_rate", "고객 규모 차이 보정"],
            ["매출 효율", "revenue_per_subscriber", "규모 대비 수익성 확인"],
            ["상호작용", "revenue_engagement_interaction", "매출과 활동 상태 결합"],
            ["변환", "log/sqrt revenue", "왜도 완화와 선형 모델 안정화"],
            ["범주형", "segment, ZIP 포함/제외/top-N", "고객군/지역 효과 검증"],
        ],
        0.65,
        1.12,
        12.05,
        3.25,
        font_size=9,
        first_col_width=2.0,
    )
    add_picture_fit(slide, ASSETS / "04_feature_importance_main.png", 1.0, 4.85, 5.5, 1.5)
    add_text(slide, "결론: 단일 컬럼보다 파생/교차 feature가 있어야 모델이 churn 신호를 더 잘 읽는다.", 6.8, 5.15, 5.4, 0.55, size=14, bold=True, color=ACCENT_2)
    add_footer(slide, page)
    page += 1

    add_section_slide(
        prs,
        page,
        "논문 재현과 비교",
        "먼저 같은 기준선을 재현한 뒤, 우리 추가 실험은 별도 기여로 해석",
        [
            "논문 best는 EasyEnsembleClassifier, F1 0.129.",
            "우리 재현 EasyEnsemble F1은 0.128로 거의 동일.",
            "LR F1 0.1681은 논문에 없던 추가 모델 조합이므로 직접 우월 주장으로 쓰지 않는다.",
        ],
    )
    page += 1

    slide = blank_slide(prs)
    add_title(slide, "논문 방식과 우리 방식", "같은 데이터셋에서 재현 가능성과 운영 질문을 동시에 확보")
    add_small_table(
        slide,
        [
            ["항목", "Makokha et al. (2026)", "우리 프로젝트"],
            ["데이터", "B2B telecom, 약 8,454 계정", "동일 원천 CSV, 로컬 8,453행"],
            ["결측 처리", "0/Unknown/median", "동일 원칙 + missing flag"],
            ["Feature", "14 raw -> 22 final 설명", "core/extended/KA/ZIP variant 분리"],
            ["Billing_ZIP", "포함 단일 설정", "포함/제외/top-N ablation"],
            ["모델", "EasyEnsemble 등 광범위 비교", "재현 + LR/Bagging/CatBoost/XGB 운영 비교"],
            ["Threshold", "F1 + recall 제약", "validation F1 + recall>=0.30, cost sweep"],
            ["Explainability", "SHAP/LIME", "LR coefficient, permutation, local contribution"],
        ],
        0.55,
        1.12,
        12.25,
        3.6,
        font_size=8,
        first_col_width=2.0,
    )
    add_text(slide, "비교 원칙: EasyEnsemble끼리는 재현 비교, LR/ZIP/cost 분석은 추가 기여로 분리한다.", 0.8, 5.55, 11.7, 0.38, size=13, bold=True, color=ACCENT)
    add_footer(slide, page)
    page += 1

    slide = blank_slide(prs)
    add_title(slide, "후보 모델과 사용 이유", "불균형 데이터에서 baseline, ensemble, categorical model을 함께 확인")
    add_small_table(
        slide,
        [
            ["모델군", "사용 이유", "최종 판단"],
            ["LogisticRegression + SVMSMOTE", "설명 가능 baseline, 선형 경계 확인", "hold-out F1 1위"],
            ["RandomForest / GradientBoosting", "비선형 관계와 상호작용 확인", "F1/recall 낮아 보조"],
            ["EasyEnsemble / RUSBoost", "논문 및 imbalance-aware 비교", "논문 재현, precision 낮음"],
            ["BalancedBagging", "minority 탐지 증가 기대", "CV 안정성과 recall 균형 좋음"],
            ["CatBoost native", "범주형 tabular 데이터 장점", "recall-heavy 후보"],
            ["XGBoost_SMOTE", "boosting 확장 후보", "recall 최고, FP 큼"],
            ["Soft ensemble", "복수 모델 평균", "결론을 바꿀 정도는 아님"],
        ],
        0.6,
        1.18,
        12.1,
        3.45,
        font_size=8,
        first_col_width=2.4,
    )
    add_bullets(slide, ["정확도 높은 모델보다 minority class를 실제로 잡는 모델을 우선 비교했다."], 0.9, 5.35, 11.5, 0.5, size=14)
    add_footer(slide, page)
    page += 1

    slide = blank_slide(prs)
    add_title(slide, "하이퍼파라미터와 Threshold 튜닝", "모델 점수보다 운영 목적에 맞는 threshold 선택이 중요")
    add_small_table(
        slide,
        [
            ["대상", "튜닝 내용", "선택 기준"],
            ["LogisticRegression", "C값 조정, class weight 비교", "validation F1"],
            ["BalancedBagging", "tree depth None/5, leaf 10/25 등", "validation threshold 후 test 적용"],
            ["CatBoost", "encoded/native categorical, depth/lr 후보", "recall-heavy 가능성 확인"],
            ["EasyEnsemble", "n=10, n=50 비교", "논문 재현 및 recall 후보"],
            ["Threshold", "0.01-0.99 또는 0.05-0.70 sweep", "F1 또는 비용 기준 net benefit"],
            ["Recall constraint", "validation recall >= 0.30", "실무형 최소 recall 확보"],
            ["Top-k", "5%-100% 고객 정렬", "예산별 캠페인 대상 수 결정"],
        ],
        0.65,
        1.15,
        12.0,
        3.35,
        font_size=9,
        first_col_width=2.2,
    )
    add_text(slide, "주의: test set에서 사후적으로 threshold를 고른 값을 최종 성능으로 쓰지 않았다.", 0.85, 5.35, 11.8, 0.35, size=13, bold=True, color=ACCENT_3)
    add_footer(slide, page)
    page += 1

    slide = blank_slide(prs)
    add_title(slide, "논문과 비교 가능한 핵심 성능", "논문 baseline은 재현했고, 우리 추가 모델은 운영 목적별 대안")
    rows = [["구분", "Model", "F1", "Recall", "Precision", "PR-AUC", "해석"]]
    for key in ["paper_reference_best", "paper_reproduction_model", "main_f1_model", "balanced_operation_model", "core_recall_heavy_model"]:
        row = final[final["selection_purpose"].eq(key)].iloc[0]
        rows.append(
            [
                key.replace("_", " ")[:18],
                short_model(row["model"]),
                fmt(row["f1"], 4),
                fmt(row["recall"], 4),
                fmt(row["precision"], 4),
                fmt(row["pr_auc"], 4),
                str(row["comment"])[:28],
            ]
        )
    add_small_table(slide, rows, 0.45, 1.1, 12.45, 2.5, font_size=7, first_col_width=1.65)
    add_picture_fit(slide, ASSETS / "05_paper_comparison.png", 0.9, 4.25, 5.0, 2.0)
    add_bullets(
        slide,
        [
            "EasyEnsemble F1 0.128은 논문 F1 0.129와 사실상 일치한다.",
            "LR F1 0.1681은 추가 발견이지만 CV에서는 평균 0.1309로 낮아진다.",
        ],
        6.45,
        4.35,
        5.9,
        1.1,
        size=13,
    )
    add_footer(slide, page)
    page += 1

    slide = blank_slide(prs)
    add_title(slide, "Hold-Out 모델 비교", "한 모델이 F1, recall, precision을 모두 이기지는 못했다")
    rows = [["목적", "Variant", "Model", "Threshold", "F1", "Recall", "Precision", "TP", "FP", "FN"]]
    for _, row in final.head(6).iterrows():
        rows.append(
            [
                short_purpose(row["selection_purpose"]),
                str(row["variant"]).replace("_", " ")[:16],
                short_model(row["model"]),
                fmt(row["threshold"], 2),
                fmt(row["f1"], 4),
                fmt(row["recall"], 4),
                fmt(row["precision"], 4),
                fmt(row["tp"], 0),
                fmt(row["fp"], 0),
                fmt(row["fn"], 0),
            ]
        )
    add_small_table(slide, rows, 0.35, 1.05, 12.65, 2.65, font_size=7, first_col_width=1.55)
    add_picture_fit(slide, ASSETS / "02_confusion_counts.png", 0.7, 4.15, 5.7, 2.35)
    add_picture_fit(slide, ASSETS / "03_precision_recall_tradeoff.png", 6.8, 4.15, 5.7, 2.35)
    add_footer(slide, page)
    page += 1

    slide = blank_slide(prs)
    add_title(slide, "5-Fold CV 안정성", "단일 hold-out 결과만으로 우위를 주장하지 않음")
    rows = [["Variant", "Model", "F1 mean", "F1 SD", "Recall", "Precision", "PR-AUC"]]
    for _, row in cv.iterrows():
        rows.append(
            [
                str(row["variant"]).replace("_", " "),
                str(row["model"]),
                fmt(row["f1_mean"], 4),
                fmt(row["f1_sd"], 4),
                fmt(row["recall_mean"], 4),
                fmt(row["precision_mean"], 4),
                fmt(row["pr_auc_mean"], 4),
            ]
        )
    add_small_table(slide, rows, 0.65, 1.18, 12.0, 2.25, font_size=8, first_col_width=2.0)
    add_bullets(
        slide,
        [
            "BalancedBagging/EasyEnsemble은 fold 평균과 표준편차가 안정적이다.",
            "LR은 hold-out F1은 높지만 CV 평균이 0.1309로 낮아져 split 민감성이 있다.",
            "논문보다 압도적 우위가 아니라 운영 목적별 대안 제시가 타당하다.",
        ],
        0.95,
        4.25,
        11.6,
        1.45,
        size=14,
    )
    add_footer(slide, page)
    page += 1

    slide = blank_slide(prs)
    add_title(slide, "Billing ZIP Ablation", "지역 변수는 모델 계열에 따라 signal과 noise가 다르게 작동")
    rows = [["Variant", "Model", "F1", "Recall", "Precision", "PR-AUC"]]
    for _, row in zip_ablation.iterrows():
        rows.append(
            [
                str(row["variant"]).replace("_", " "),
                str(row["model"]),
                fmt(row["f1"], 4),
                fmt(row["recall"], 4),
                fmt(row["precision"], 4),
                fmt(row["pr_auc"], 4),
            ]
        )
    add_small_table(slide, rows, 0.55, 1.18, 12.25, 2.7, font_size=8, first_col_width=2.1)
    add_bullets(
        slide,
        [
            "BalancedBagging은 ZIP 포함이 F1/recall을 높였다.",
            "Logistic Regression은 ZIP 제외가 F1 0.1681로 가장 좋았다.",
            "고카디널리티 label encoding은 선형 모델에서 noise가 될 수 있다.",
        ],
        0.95,
        4.65,
        11.4,
        1.15,
        size=14,
    )
    add_footer(slide, page)
    page += 1

    slide = blank_slide(prs)
    add_title(slide, "Feature Group Ablation", "어떤 feature 그룹이 어느 모델에서 중요한지 정량화")
    subset = feature_ablation[
        (feature_ablation["experiment"].eq("DROP_GROUP"))
        & (
            ((feature_ablation["model"].eq("LogisticRegression_SMOTE")) & (feature_ablation["variant"].eq("without_billing_zip")))
            | ((feature_ablation["model"].eq("BalancedBagging_original")) & (feature_ablation["variant"].eq("with_billing_zip")))
        )
    ].copy()
    subset = subset.sort_values("f1_delta_vs_baseline").head(8)
    rows = [["Model", "Dropped group", "F1", "Baseline", "Delta"]]
    for _, row in subset.iterrows():
        rows.append(
            [
                "LR" if "Logistic" in row["model"] else "BalancedBagging",
                row["group"],
                fmt(row["f1"], 4),
                fmt(row["baseline_f1"], 4),
                fmt(row["f1_delta_vs_baseline"], 4),
            ]
        )
    add_small_table(slide, rows, 0.65, 1.15, 7.5, 3.3, font_size=8, first_col_width=1.45)
    add_bullets(
        slide,
        [
            "LR: categorical group 제거 시 F1이 0.1681에서 0.0806으로 급락.",
            "BalancedBagging: interaction group 제거 시 가장 큰 성능 하락.",
            "feature engineering 효과는 모델 구조와 결합해 해석해야 한다.",
        ],
        8.45,
        1.55,
        4.1,
        2.2,
        size=13,
    )
    add_footer(slide, page)
    page += 1

    add_section_slide(
        prs,
        page,
        "컬럼별 데이터 분리와 단일 컬럼 실험",
        "각 컬럼만 넣었을 때 무엇이 좋고 나쁜지 확인해 최종 feature 설계를 방어",
        [
            "원본 CSV는 수정하지 않고 컬럼별 `feature + CHURN` CSV를 생성했다.",
            "범주형 값별 하위 데이터와 Yes/No 이탈률 summary를 만들었다.",
            "단일 컬럼 모델은 성능 상한을 확인하는 빠른 검증용이다.",
        ],
    )
    page += 1

    slide = blank_slide(prs)
    add_title(slide, "컬럼별 CSV 생성 결과", "각 변수의 역할과 이탈률을 따로 설명하기 위한 보조 데이터")
    add_small_table(
        slide,
        [
            ["생성 위치", "내용", "발표에서 쓰는 이유"],
            ["01_column_churn_pairs", "각 원본 feature + CHURN", "단일 변수 설명"],
            ["02_category_value_subsets", "Bronze, SOHO, ZIP별 하위 CSV", "범주 값별 이탈률 확인"],
            ["03_profiles", "전체/범주/수치 요약", "EDA와 보고서 표 근거"],
            ["05_single_column_model_screening.csv", "단일 컬럼 모델 전체 결과", "컬럼별 성능 한계 확인"],
            ["05_single_column_best_models.csv", "컬럼별 best model", "좋은/나쁜 컬럼 순위"],
            ["00_data_dictionary_korean.md", "한글 컬럼 사전", "발표 이해도 보강"],
        ],
        0.75,
        1.18,
        11.85,
        2.75,
        font_size=9,
        first_col_width=3.1,
    )
    add_bullets(
        slide,
        [
            "PID와 KA_name은 모델용/분리용 데이터에서는 제외했다.",
            "Billing_ZIP은 고유값이 많아 값별 이탈률과 모델별 효과를 별도로 확인했다.",
        ],
        1.0,
        4.85,
        11.2,
        1.0,
        size=14,
    )
    add_footer(slide, page)
    page += 1

    slide = blank_slide(prs)
    add_title(slide, "단일 컬럼 모델 스크리닝", "컬럼 하나만 넣으면 F1은 약 0.15 근처에서 한계")
    rows = [["Column", "Type", "Best model", "F1", "Recall", "Precision", "PR-AUC"]]
    for _, row in single.head(10).iterrows():
        rows.append(
            [
                row["column"],
                row["semantic_type"],
                row["model"].replace("_balanced", ""),
                fmt(row["f1"], 4),
                fmt(row["recall"], 4),
                fmt(row["precision"], 4),
                fmt(row["pr_auc"], 4),
            ]
        )
    add_small_table(slide, rows, 0.5, 1.05, 12.35, 4.3, font_size=7, first_col_width=2.0)
    add_text(slide, "해석: 가입자 수와 매출 규모가 상대적으로 강하지만, 단일 변수만으로는 충분하지 않다.", 0.8, 5.95, 11.8, 0.35, size=13, bold=True, color=ACCENT)
    add_footer(slide, page)
    page += 1

    slide = blank_slide(prs)
    add_title(slide, "컬럼별로 좋았던 점과 안 좋았던 점", "단일 컬럼 결과를 최종 feature 설계로 연결")
    add_small_table(
        slide,
        [
            ["컬럼군", "좋았던 점", "안 좋았던 점", "최종 반영"],
            ["가입자 수", "Total_SUBs F1 0.1498, 활동 규모 신호", "규모 자체만으로는 precision 낮음", "active/inactive/dormant rate"],
            ["모바일/총매출", "AvgMobileRevenue, TotalRevenue 상위", "왜도와 outlier 큼", "log/sqrt, revenue ratio"],
            ["EffectiveSegment", "고객군별 차이 존재", "단독 F1 0.1346로 제한", "categorical/frequency encoding"],
            ["CRM value", "recall은 높게 나올 수 있음", "precision 낮음", "segment + interaction"],
            ["Billing_ZIP", "단독 recall 0.7000", "고카디널리티와 과적합 위험", "포함/제외/top-N ablation"],
            ["Suspended", "결측/존재 여부 신호 가능", "결측률 높아 단독 모델은 약함", "missing flag + 0 대체"],
        ],
        0.4,
        1.05,
        12.55,
        4.2,
        font_size=7,
        first_col_width=1.7,
    )
    add_bullets(slide, ["결론: 단일 컬럼 순위는 최종 모델 순위가 아니라 feature engineering 방향을 정하는 근거다."], 0.75, 5.85, 11.8, 0.45, size=13)
    add_footer(slide, page)
    page += 1

    slide = blank_slide(prs)
    add_title(slide, "Paper/KA Ablation Variants", "논문형 core, ZIP top-N, KA 추상화 feature를 실제 성능으로 점검")
    rows = [["Variant", "Feature", "Model", "Threshold", "F1", "Recall", "Precision", "PR-AUC"]]
    for _, row in paper_ablation.head(6).iterrows():
        rows.append(
            [
                str(row["variant"])[:32],
                fmt(row["feature_count"], 0),
                str(row["model"]).replace("BalancedBagging_", "BB_")[:22],
                fmt(row["threshold"], 2),
                fmt(row["f1"], 4),
                fmt(row["recall"], 4),
                fmt(row["precision"], 4),
                fmt(row["pr_auc"], 4),
            ]
        )
    add_small_table(slide, rows, 0.35, 1.05, 12.65, 2.75, font_size=6, first_col_width=3.1)
    add_bullets(
        slide,
        [
            "최고 F1은 `paper_core_zip_log_ka_abstract + BalancedBagging`의 0.1561.",
            "기존 최종 LR F1 0.1681을 넘지는 못했다.",
            "KA 실명 대신 추상 feature를 쓰는 방식은 발표 보조 근거로 활용 가능하다.",
        ],
        0.9,
        4.45,
        11.4,
        1.35,
        size=13,
    )
    add_footer(slide, page)
    page += 1

    slide = blank_slide(prs)
    add_title(slide, "추가 모델/하이퍼파라미터 실험", "BalancedBagging 튜닝은 좋아졌지만 최종 결론은 바꾸지 않음")
    rows = [["Variant", "Model", "Th.", "F1", "Recall", "Precision", "TP", "FP"]]
    for _, row in additional.head(8).iterrows():
        rows.append(
            [
                str(row["variant"]).replace("_", " ")[:18],
                str(row["model"])[:28],
                fmt(row["selected_threshold"], 2),
                fmt(row["test_f1"], 4),
                fmt(row["test_recall"], 4),
                fmt(row["test_precision"], 4),
                fmt(row["test_tp"], 0),
                fmt(row["test_fp"], 0),
            ]
        )
    add_small_table(slide, rows, 0.45, 1.05, 12.45, 3.15, font_size=7, first_col_width=1.8)
    add_bullets(
        slide,
        [
            "추가 실험 best: `BalancedBagging_tree_depthnone_leaf25`, test F1 0.1605.",
            "기존 최종 모델 `LogisticRegression_SMOTE` F1 0.1681은 유지.",
            "튜닝 결과는 recall 운영 후보를 설명하는 보조 근거로 사용한다.",
        ],
        0.9,
        4.95,
        11.5,
        1.25,
        size=13,
    )
    add_footer(slide, page)
    page += 1

    slide = blank_slide(prs)
    add_title(slide, "CRM Segment 분석", "전체 평균만 보면 고객군별 실패 패턴을 놓친다")
    rows = [["Segment", "Rows", "Positives", "Recall", "Precision", "F1", "FN revenue risk"]]
    for _, row in segment.iterrows():
        rows.append(
            [
                row["CRM_segment_bucket"],
                fmt(row["rows"], 0),
                fmt(row["positives"], 0),
                fmt(row["recall"], 4),
                fmt(row["precision"], 4),
                fmt(row["f1"], 4),
                fmt(row["fn_total_revenue_at_risk"], 2),
            ]
        )
    add_small_table(slide, rows, 0.75, 1.15, 11.85, 1.7, font_size=9, first_col_width=2.1)
    add_picture_fit(slide, PHASE6 / "phase6_segment_recall_heatmap.png", 0.85, 3.35, 6.1, 2.6)
    add_bullets(
        slide,
        [
            "mid/high value는 recall이 높지만 precision이 낮아 FP 관리가 중요하다.",
            "low value는 놓치는 이탈자가 많아 FN 위험이 크다.",
        ],
        7.35,
        3.65,
        4.9,
        1.2,
        size=13,
    )
    add_footer(slide, page)
    page += 1

    slide = blank_slide(prs)
    add_title(slide, "High-Value 전용 모델 실험", "전용 모델은 precision을 높였지만 recall을 크게 낮췄다")
    rows = [["Scope", "Train", "Test", "F1", "Recall", "Precision", "PR-AUC", "FP", "FN"]]
    for _, row in high_value.iterrows():
        rows.append(
            [
                row["model_scope"].replace("_", " ")[:28],
                fmt(row["train_rows"], 0),
                fmt(row["test_rows"], 0),
                fmt(row["f1"], 4),
                fmt(row["recall"], 4),
                fmt(row["precision"], 4),
                fmt(row["pr_auc"], 4),
                fmt(row["fp"], 0),
                fmt(row["fn"], 0),
            ]
        )
    add_small_table(slide, rows, 0.65, 1.25, 12.05, 1.55, font_size=8, first_col_width=2.8)
    add_bullets(
        slide,
        [
            "Global high-recall 모델은 high-value에서도 recall 0.7857이지만 FP가 314.",
            "High-value-only 모델은 precision/PR-AUC가 개선되지만 recall 0.3571로 낮아진다.",
            "따라서 전용 모델 단독 교체보다 two-stage review가 적합하다.",
        ],
        0.95,
        3.75,
        11.3,
        1.55,
        size=14,
    )
    add_footer(slide, page)
    page += 1

    slide = blank_slide(prs)
    add_title(slide, "Cost-Sensitive Threshold", "논문 비용 구조에서는 recall 극대화가 순이익에 유리")
    rows = [["Scenario", "Cost ratio", "Threshold", "Net value", "Recall", "Precision", "F1"]]
    for _, row in cost_best.iterrows():
        rows.append(
            [
                row["scenario"],
                fmt(row["cost_ratio"], 1),
                fmt(row["threshold"], 2),
                fmt(row["expected_net_value"], 0),
                fmt(row["recall"], 4),
                fmt(row["precision"], 4),
                fmt(row["f1"], 4),
            ]
        )
    add_small_table(slide, rows, 0.55, 1.12, 12.25, 2.65, font_size=8, first_col_width=2.3)
    add_picture_fit(slide, PHASE6 / "phase6_cost_best_paper_baseline.png", 1.0, 4.25, 5.2, 1.9)
    add_bullets(
        slide,
        [
            "FP cost가 커지는 보수적 시나리오에서는 기대가치가 음수로 전환될 수 있다.",
            "threshold는 확률값이 아니라 운영 sweep 결과로 해석해야 한다.",
        ],
        6.7,
        4.45,
        5.6,
        1.05,
        size=13,
    )
    add_footer(slide, page)
    page += 1

    slide = blank_slide(prs)
    add_title(slide, "비즈니스 임팩트 시나리오", "모델 선택은 팀 역량과 캠페인 비용에 따라 바뀐다")
    add_picture_fit(slide, PHASE5B / "business_impact_dashboard.png", 0.55, 1.05, 12.25, 5.65)
    add_footer(slide, page)
    page += 1

    slide = blank_slide(prs)
    add_title(slide, "Top-k 예산 전략", "실제 캠페인은 상위 몇 %를 접촉할지로 운영하기 쉽다")
    best_topk = (
        topk.sort_values(["topk_pct", "net_benefit"], ascending=[True, False])
        .groupby("topk_pct")
        .head(1)
        .query("topk_pct in [0.05, 0.1, 0.2, 0.3, 0.4, 0.5]")
    )
    rows = [["Top-k", "Best case", "Contacts", "TP", "Recall@k", "Precision@k", "Net benefit"]]
    for _, row in best_topk.iterrows():
        rows.append(
            [
                pct(row["topk_pct"], 0),
                row["case_label"][:25],
                fmt(row["contacts"], 0),
                fmt(row["tp"], 0),
                fmt(row["recall_at_k"], 4),
                fmt(row["precision_at_k"], 4),
                fmt(row["net_benefit"], 0),
            ]
        )
    add_small_table(slide, rows, 0.55, 1.12, 12.25, 2.15, font_size=8, first_col_width=1.2)
    add_picture_fit(slide, PHASE6 / "phase6_topk_budget_curves.png", 1.1, 3.65, 5.8, 2.6)
    add_bullets(
        slide,
        [
            "예산이 작으면 LR/EasyEnsemble이 효율적이다.",
            "예산이 커질수록 BalancedBagging/CatBoost 계열의 recall 장점이 커진다.",
        ],
        7.35,
        3.95,
        4.9,
        1.1,
        size=13,
    )
    add_footer(slide, page)
    page += 1

    slide = blank_slide(prs)
    add_title(slide, "Calibration", "Raw score는 실제 이탈 확률이 아니므로 보정이 필요")
    cal_rows = [["Case", "Method", "Brier", "ECE", "Mean score", "Observed"]]
    cal_show = calibration[
        (calibration["case_label"].isin(["LR_no_zip_f1", "BalancedBagging_with_zip", "CatBoost_native_with_zip", "XGBoost_with_zip"]))
        & (calibration["calibration_method"].isin(["raw", "platt"]))
    ]
    for _, row in cal_show.iterrows():
        cal_rows.append(
            [
                short_case(row["case_label"]),
                row["calibration_method"],
                fmt(row["brier_score"], 4),
                fmt(row["ece_10bin"], 4),
                fmt(row["mean_score"], 4),
                fmt(row["observed_churn_rate"], 4),
            ]
        )
    add_small_table(slide, cal_rows, 0.55, 1.1, 12.25, 3.0, font_size=7, first_col_width=2.4)
    add_picture_fit(slide, PHASE6 / "phase6_calibration_comparison.png", 1.0, 4.55, 5.5, 1.7)
    add_text(slide, "발표 문장: score 0.4를 이탈 확률 40%라고 말하면 안 된다.", 6.9, 5.05, 5.6, 0.45, size=13, bold=True, color=ACCENT_3)
    add_footer(slide, page)
    page += 1

    slide = blank_slide(prs)
    add_title(slide, "Feature Importance와 해석 가능성", "논문 SHAP과 우리 LR 해석은 순위는 달라도 결론은 수렴")
    add_picture_fit(slide, ROOT / "processed" / "phase_5a_interpretability" / "lr_coefficient_importance.png", 0.65, 1.2, 5.7, 3.2)
    add_picture_fit(slide, ROOT / "processed" / "phase_5a_interpretability" / "lr_linear_contribution_importance.png", 6.85, 1.2, 5.7, 3.2)
    add_small_table(
        slide,
        [
            ["해석 방법", "상위 신호", "의미"],
            ["논문 SHAP", "active subscriber rate", "tree ensemble의 비선형 활동성 효과"],
            ["Permutation FI", "AvgMobileRevenue_sqrt, TotalRevenue_sqrt", "교란 시 F1 손실이 큰 feature"],
            ["LR coefficient", "AvgFIXRevenue_log, AvgMobileRevenue_sqrt", "log-odds 방향성과 크기"],
            ["Local contribution", "coefficient x standardized value", "고객별 예측 근거"],
        ],
        0.75,
        4.85,
        11.85,
        1.35,
        font_size=8,
        first_col_width=2.1,
    )
    add_footer(slide, page)
    page += 1

    slide = blank_slide(prs)
    add_title(slide, "모델 합의도 기반 Risk Tier", "여러 모델이 동시에 경고한 고객군은 실제 이탈률이 높다")
    add_picture_fit(slide, PHASE6 / "phase6_model_agreement.png", 0.85, 1.2, 6.1, 4.0)
    add_small_table(
        slide,
        [
            ["Vote count", "고객 수", "실제 이탈자", "관측 이탈률"],
            ["0", "91", "5", "5.49%"],
            ["1", "195", "6", "3.08%"],
            ["6", "355", "22", "6.20%"],
            ["7", "134", "13", "9.70%"],
            ["8", "169", "21", "12.43%"],
        ],
        7.35,
        1.35,
        4.8,
        1.9,
        font_size=9,
        first_col_width=1.35,
    )
    add_bullets(
        slide,
        [
            "8개 모델 모두가 경고한 고객군은 전체 평균 이탈률 6.46%의 약 1.9배.",
            "합의도는 영업팀 우선순위 tier로 쓸 수 있다.",
            "단, vote 0에서도 이탈자가 있어 완전한 rule은 아니다.",
        ],
        7.45,
        3.75,
        4.8,
        1.45,
        size=13,
    )
    add_footer(slide, page)
    page += 1

    slide = blank_slide(prs)
    add_title(slide, "Bootstrap CI와 McNemar 검정", "모델 차이는 단순 점수 차이가 아니라 오류 패턴 차이")
    rows = [["Case", "F1 point", "F1 95% CI", "Recall point", "Recall 95% CI", "Precision point"]]
    for case in ["LR_no_zip_f1", "BalancedBagging_with_zip", "CatBoost_native_with_zip", "XGBoost_with_zip"]:
        row = metric_ci[metric_ci["case_label"].eq(case)].iloc[0]
        rows.append(
            [
                short_case(case),
                fmt(row["f1_point"], 4),
                f"[{fmt(row['f1_ci95_low'], 3)}, {fmt(row['f1_ci95_high'], 3)}]",
                fmt(row["recall_point"], 4),
                f"[{fmt(row['recall_ci95_low'], 3)}, {fmt(row['recall_ci95_high'], 3)}]",
                fmt(row["precision_point"], 4),
            ]
        )
    add_small_table(slide, rows, 0.55, 1.15, 12.25, 2.15, font_size=8, first_col_width=2.4)
    add_bullets(
        slide,
        [
            "LR의 F1 point estimate는 높지만 신뢰구간이 넓다.",
            "BalancedBagging/CatBoost/XGBoost는 recall 중심 운영 장점이 일관된다.",
            "McNemar test는 주요 모델의 정오분류 패턴이 유의하게 다름을 보여준다.",
        ],
        0.95,
        4.2,
        11.2,
        1.35,
        size=14,
    )
    add_footer(slide, page)
    page += 1

    slide = blank_slide(prs)
    add_title(slide, "재현 자동화와 MLOps", "제출 이후 운영형 프로젝트로 확장 가능한 구조")
    add_process_lane(
        slide,
        [
            ("pytest", "데이터\n스키마 검사"),
            ("preprocess", "전처리/모델\n평가 생성"),
            ("phase 3-8", "실험/검정\n재실행"),
            ("PPT", "이미지와\nPPT 재생성"),
            ("n8n", "Docker runner\nworkflow"),
        ],
        0.85,
        1.35,
        11.5,
    )
    add_small_table(
        slide,
        [
            ["구성", "구현 내용"],
            ["n8n workflow", "Health -> Full Reproduction -> Summary"],
            ["Runner API", "X-API-KEY 헤더 인증"],
            ["Drift check", "monitor_drift.py의 PSI 기반 점검"],
            ["재학습 기준", "F1 0.14 이하 또는 PSI 0.2 초과 시 검토"],
            ["버전 관리", "대량 processed 산출물은 DVC/외부 스토리지 권장"],
        ],
        1.0,
        3.7,
        11.2,
        2.0,
        font_size=9,
        first_col_width=2.3,
    )
    add_footer(slide, page)
    page += 1

    slide = blank_slide(prs)
    add_title(slide, "한계와 향후 개선", "성능 상한의 핵심은 모델보다 데이터 구조")
    add_small_table(
        slide,
        [
            ["한계", "현재 영향", "개선 방향"],
            ["정적 CRM snapshot", "이탈 직전 행동 변화 없음", "월별 사용량/매출 추세 추가"],
            ["결제/VOC/계약 정보 없음", "이탈 원인 직접 신호 부족", "결제 실패, 문의, 계약 만료 feature"],
            ["외부 검증 없음", "generalization 확인 제한", "기간별 hold-out, 다른 시장 데이터"],
            ["비용 가정 고정", "운영 결론이 비용에 민감", "고객별 ARPU 기반 individualized value"],
            ["확률 보정 미배포", "score를 확률로 쓰기 어려움", "calibrated probability threshold"],
        ],
        0.65,
        1.15,
        12.0,
        2.65,
        font_size=8,
        first_col_width=2.2,
    )
    add_bullets(
        slide,
        [
            "모델을 더 추가하는 것보다 temporal feature 확보가 성능 개선의 핵심이다.",
            "낮은 F1은 실험 실패가 아니라 강한 불균형과 정적 데이터 한계의 결과로 설명한다.",
        ],
        0.95,
        4.6,
        11.4,
        1.1,
        size=14,
    )
    add_footer(slide, page)
    page += 1

    slide = blank_slide(prs)
    add_band(slide, 0, 0, 13.333, 7.5, RGBColor(235, 242, 248))
    add_text(slide, "최종 결론", 0.8, 0.8, 11.8, 0.7, size=36, bold=True, color=TITLE)
    add_bullets(
        slide,
        [
            "논문 EasyEnsemble baseline은 F1 0.128로 재현했다.",
            "Hold-out F1 최고는 LogisticRegression_SMOTE지만, CV 안정성은 BalancedBagging/EasyEnsemble이 강하다.",
            "컬럼별 실험은 가입자 수/매출/segment/ZIP이 단독으로는 약하고, 파생/교차 feature가 필요함을 보여줬다.",
            "비용, 예산, recall 목표, segment에 따라 최적 모델과 threshold는 달라진다.",
            "따라서 ChurnRadar의 답은 단일 최고 모델이 아니라 운영 목적별 모델 선택 프레임워크다.",
        ],
        1.0,
        2.0,
        11.0,
        3.5,
        size=16,
    )
    add_text(slide, "마지막 발표 문장: 높은 점수 하나보다, 어떤 조건에서 어떤 모델을 써야 하는지 설명하는 것이 이 프로젝트의 핵심입니다.", 1.0, 6.35, 11.3, 0.45, size=13, bold=True, color=ACCENT)
    add_footer(slide, page)

    return prs


def main() -> None:
    prs = build_detailed_presentation()
    prs.save(OUTPUT)
    if SUBMISSION_OUTPUT.parent.exists():
        prs.save(SUBMISSION_OUTPUT)
    print(f"Saved: {OUTPUT}")
    if SUBMISSION_OUTPUT.exists():
        print(f"Saved: {SUBMISSION_OUTPUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
