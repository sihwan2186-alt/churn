from pathlib import Path
from typing import Iterable

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUTPUT = ROOT / "ChurnRadar_Final_Presentation.pptx"
ASSETS = ROOT / "presentation_assets"
PHASE5B = ROOT / "processed" / "phase_5b_business_impact"
PHASE6 = ROOT / "processed" / "phase_6_extended_case_studies"
PHASE8 = ROOT / "processed" / "phase_8_statistical_validation"

FONT = "Malgun Gothic"
TITLE = RGBColor(22, 38, 58)
TEXT = RGBColor(44, 52, 64)
MUTED = RGBColor(97, 108, 121)
ACCENT = RGBColor(31, 119, 180)
ACCENT_2 = RGBColor(46, 160, 120)
ACCENT_3 = RGBColor(222, 125, 44)
LIGHT = RGBColor(244, 247, 250)
WHITE = RGBColor(255, 255, 255)


def set_run(run, size: int, bold: bool = False, color: RGBColor = TEXT) -> None:
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_text(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    size: int = 18,
    bold: bool = False,
    color: RGBColor = TEXT,
    align=PP_ALIGN.LEFT,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color)


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    add_text(slide, title, 0.55, 0.25, 12.2, 0.45, size=23, bold=True, color=TITLE)
    if subtitle:
        add_text(slide, subtitle, 0.58, 0.74, 11.8, 0.35, size=11, color=MUTED)


def add_footer(slide, page: int) -> None:
    add_text(slide, f"ChurnRadar | {page}", 11.55, 7.08, 1.25, 0.22, size=8, color=MUTED, align=PP_ALIGN.RIGHT)


def add_band(slide, left: float, top: float, width: float, height: float, fill: RGBColor = LIGHT) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill


def add_callout(slide, text: str, left: float, top: float, width: float, height: float, color: RGBColor = ACCENT) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = Inches(0.12)
    frame.margin_right = Inches(0.12)
    frame.margin_top = Inches(0.08)
    frame.margin_bottom = Inches(0.08)
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    set_run(run, size=15, bold=True, color=WHITE)


def add_bullets(
    slide,
    bullets: Iterable[str],
    left: float,
    top: float,
    width: float,
    height: float,
    size: int = 14,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for idx, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.space_after = Pt(4)
        paragraph.font.name = FONT
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = TEXT


def add_picture_fit(slide, path: Path, left: float, top: float, width: float, height: float) -> None:
    if not path.exists():
        add_band(slide, left, top, width, height, RGBColor(250, 238, 238))
        add_text(slide, f"Missing image:\n{path.name}", left + 0.15, top + 0.2, width - 0.3, height - 0.4, size=12, color=RGBColor(140, 40, 40))
        return
    slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width), height=Inches(height))


def add_table(slide, rows: list[list[str]], left: float, top: float, width: float, height: float, header_fill: RGBColor = TITLE) -> None:
    row_count = len(rows)
    col_count = len(rows[0])
    table_shape = slide.shapes.add_table(
        row_count,
        col_count,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    table = table_shape.table
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(value)
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    set_run(
                        run,
                        size=9 if r else 9,
                        bold=(r == 0),
                        color=WHITE if r == 0 else TEXT,
                    )
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(248, 250, 252)


def blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = blank_slide(prs)
    add_band(slide, 0, 0, 13.333, 7.5, RGBColor(236, 242, 248))
    add_text(slide, "ChurnRadar", 0.75, 1.0, 12.0, 0.8, size=44, bold=True, color=TITLE)
    add_text(slide, "B2B 통신사 고객 이탈 예측과 운영 시나리오 분석", 0.82, 1.85, 11.0, 0.45, size=20, color=TEXT)
    add_callout(slide, "F1 | Recall | Precision | Cost", 0.82, 2.62, 4.6, 0.55, ACCENT)
    add_bullets(
        slide,
        [
            "이탈 비율 6.5%의 불균형 데이터",
            "논문 EasyEnsemble baseline 재현",
            "운영 목적별 모델 선택 프레임워크 제시",
        ],
        0.86,
        3.55,
        6.0,
        1.5,
        size=16,
    )
    add_picture_fit(slide, ASSETS / "01_model_metric_comparison.png", 6.6, 2.05, 5.9, 3.5)
    add_footer(slide, 1)

    slide = blank_slide(prs)
    add_title(slide, "데이터와 문제 정의", "정확도보다 minority class 탐지가 중요한 churn prediction 문제")
    add_table(
        slide,
        [
            ["항목", "값"],
            ["원본 데이터", "8,453 rows x 14 columns"],
            ["PID 중복 제거 후", "8,436 rows"],
            ["원본 CHURN=Yes", "549"],
            ["원본 CHURN=No", "7,904"],
            ["이탈 비율", "6.49%"],
        ],
        0.8,
        1.35,
        4.7,
        2.3,
    )
    add_bullets(
        slide,
        [
            "Accuracy만 보면 대부분 고객을 비이탈로 예측해도 좋아 보인다.",
            "주요 지표는 F1, recall, precision, PR-AUC, MCC다.",
            "Test set에는 resampling을 적용하지 않아 실제 운영 분포를 유지했다.",
        ],
        6.0,
        1.35,
        6.3,
        2.1,
        size=15,
    )
    add_callout(slide, "핵심: 적은 이탈자를 놓치지 않되, FP 비용도 같이 관리", 1.0, 4.65, 11.2, 0.65, ACCENT_2)
    add_footer(slide, 2)

    slide = blank_slide(prs)
    add_title(slide, "전처리와 Leakage 방지", "모델 점수보다 먼저 재현 가능하고 안전한 split 원칙 확보")
    add_bullets(
        slide,
        [
            "PID 기준 중복 제거",
            "CHURN 이진 변환 후 feature matrix에서 제거",
            "80:20 stratified train/test split",
            "Imputation, encoding, scaling은 train 기준 fit",
            "SVMSMOTE는 train partition에만 적용",
            "Threshold는 validation에서 선택 후 test에 1회 적용",
        ],
        0.75,
        1.25,
        5.8,
        4.3,
        size=15,
    )
    add_band(slide, 6.9, 1.3, 5.35, 4.2, RGBColor(238, 246, 250))
    add_text(slide, "Leakage Guardrails", 7.25, 1.65, 4.7, 0.35, size=20, bold=True, color=TITLE)
    add_text(
        slide,
        "Train fold에서만 학습한 전처리 artifact를 validation/test에 적용했다. 따라서 test row 정보가 train preprocessing에 섞이지 않는다.",
        7.25,
        2.25,
        4.6,
        1.35,
        size=14,
    )
    add_text(
        slide,
        "불균형 처리는 성능 과장을 막기 위해 test set에는 적용하지 않았다.",
        7.25,
        3.75,
        4.6,
        0.8,
        size=14,
        color=MUTED,
    )
    add_footer(slide, 3)

    slide = blank_slide(prs)
    add_title(slide, "논문 Baseline 재현", "EasyEnsemble 기준 F1이 논문 보고값과 거의 일치")
    add_table(
        slide,
        [
            ["기준", "Model", "F1", "Recall", "Precision", "PR-AUC"],
            ["Makokha et al.", "EasyEnsemble", "0.129", "0.382", "0.077", "0.079"],
            ["우리 재현", "EasyEnsemble + ZIP", "0.128", "0.587", "0.072", "0.085"],
        ],
        0.8,
        1.25,
        11.8,
        1.25,
    )
    add_picture_fit(slide, ASSETS / "05_paper_comparison.png", 1.25, 3.0, 5.1, 2.9)
    add_bullets(
        slide,
        [
            "직접 비교는 논문 best 모델인 EasyEnsemble끼리 수행한다.",
            "LR F1 0.1681은 논문에 없는 추가 발견이므로 별도 해석한다.",
            "이후 분석은 재현 baseline 위의 운영적 확장이다.",
        ],
        7.0,
        3.05,
        5.4,
        2.2,
        size=14,
    )
    add_footer(slide, 4)

    slide = blank_slide(prs)
    add_title(slide, "Hold-Out 모델 비교", "F1, recall, precision은 같은 방향으로 움직이지 않는다")
    add_table(
        slide,
        [
            ["목적", "Model", "F1", "Recall", "Precision", "TP", "FP"],
            ["F1", "LR SMOTE", "0.1681", "0.2661", "0.1229", "29", "207"],
            ["균형", "BalancedBagging", "0.1526", "0.5872", "0.0877", "64", "666"],
            ["Recall", "CatBoost native", "0.1310", "0.8349", "0.0711", "91", "1,189"],
            ["Recall 확장", "XGBoost", "0.1242", "0.9266", "0.0665", "101", "1,417"],
        ],
        0.55,
        1.15,
        12.25,
        1.9,
    )
    add_picture_fit(slide, ASSETS / "02_confusion_counts.png", 0.8, 3.55, 5.6, 2.8)
    add_picture_fit(slide, ASSETS / "03_precision_recall_tradeoff.png", 6.8, 3.55, 5.6, 2.8)
    add_footer(slide, 5)

    slide = blank_slide(prs)
    add_title(slide, "Cross-Validation 안정성", "단일 split 최고점과 fold 평균은 다르게 해석해야 한다")
    add_table(
        slide,
        [
            ["Model", "CV F1 mean", "CV F1 SD", "Recall", "Precision"],
            ["BalancedBagging + ZIP", "0.1455", "0.0126", "0.5248", "0.0845"],
            ["EasyEnsemble + ZIP", "0.1445", "0.0117", "0.5835", "0.0824"],
            ["EasyEnsemble no ZIP", "0.1408", "0.0081", "0.5835", "0.0801"],
            ["LR no ZIP", "0.1309", "0.0154", "0.1743", "0.1053"],
        ],
        0.8,
        1.2,
        7.0,
        2.0,
    )
    add_bullets(
        slide,
        [
            "LR hold-out F1 0.1681은 좋은 단일 split 결과다.",
            "5-fold CV에서는 BalancedBagging/EasyEnsemble 계열이 더 안정적이다.",
            "결론은 특정 모델의 압도적 우위가 아니라 목적별 선택이다.",
        ],
        8.2,
        1.25,
        4.3,
        2.0,
        size=14,
    )
    add_callout(slide, "안전한 주장: 논문 baseline 재현 + 운영 목적별 대안 제시", 1.15, 4.6, 11.0, 0.7, ACCENT)
    add_footer(slide, 6)

    slide = blank_slide(prs)
    add_title(slide, "차별화 실험", "논문이 단일 설정으로 둔 요소를 분리해 확인")
    add_table(
        slide,
        [
            ["실험", "핵심 결과"],
            ["ZIP ablation", "ZIP은 tree ensemble에는 도움, LR에는 noise 가능성"],
            ["Feature group", "LR은 categorical, BalancedBagging은 interaction에 민감"],
            ["Segment", "high-value는 FP 문제, low-value는 FN 문제"],
            ["Cost threshold", "낮은 threshold는 순이익을 올리지만 접촉 수 급증"],
            ["Phase 7", "추가 후보는 결론을 바꿀 정도는 아님"],
        ],
        0.85,
        1.15,
        11.55,
        2.25,
    )
    add_picture_fit(slide, ASSETS / "04_feature_importance_main.png", 1.1, 3.9, 5.4, 2.6)
    add_bullets(
        slide,
        [
            "추가 실험의 가치는 점수 상승보다 해석력에 있다.",
            "어떤 feature가 어떤 모델에서 중요한지 정량화했다.",
        ],
        7.1,
        4.1,
        5.1,
        1.3,
        size=14,
    )
    add_footer(slide, 7)

    slide = blank_slide(prs)
    add_title(slide, "Feature Importance와 해석 가능성", "LR은 logit space에서 additive explanation 제공")
    add_bullets(
        slide,
        [
            "논문 SHAP 1위: active subscriber rate",
            "우리 LR 중요도 상위: revenue transform + engagement interaction",
            "모델 구조와 중요도 측정 방식이 달라 순위는 달라질 수 있다.",
            "큰 결론은 동일하다: 매출 패턴과 서비스 참여도 결합이 중요하다.",
        ],
        0.8,
        1.25,
        5.5,
        3.1,
        size=14,
    )
    add_picture_fit(slide, ROOT / "processed" / "phase_5a_interpretability" / "lr_coefficient_importance.png", 6.65, 1.35, 5.7, 4.2)
    add_footer(slide, 8)

    slide = blank_slide(prs)
    add_title(slide, "비즈니스 임팩트", "모델 성능은 캠페인 비용과 처리 역량에 따라 다르게 해석")
    add_picture_fit(slide, PHASE5B / "business_impact_dashboard.png", 0.65, 1.05, 12.0, 5.7)
    add_footer(slide, 9)

    slide = blank_slide(prs)
    add_title(slide, "Top-k 캠페인 예산 전략", "운영팀에는 threshold보다 상위 k% 방식이 직관적")
    add_picture_fit(slide, PHASE6 / "phase6_topk_budget_curves.png", 0.75, 1.15, 7.1, 5.0)
    add_bullets(
        slide,
        [
            "Top 10%에서는 LR이 효율적이다.",
            "Top 40% 이상에서는 BalancedBagging 계열이 유리해진다.",
            "Top 100%는 수치상 이익이어도 실제 운영 전략으로는 부적절하다.",
        ],
        8.25,
        1.55,
        4.4,
        2.1,
        size=14,
    )
    add_footer(slide, 10)

    slide = blank_slide(prs)
    add_title(slide, "비용 시나리오별 Threshold 민감도", "비용 구조가 바뀌면 최적 모델도 바뀐다")
    add_picture_fit(slide, PHASE6 / "phase6_cost_best_paper_baseline.png", 0.8, 1.25, 6.0, 4.4)
    add_bullets(
        slide,
        [
            "논문 비용 구조에서는 FN 비용이 FP 비용보다 훨씬 크다.",
            "따라서 낮은 threshold와 recall 극대화가 유리하다.",
            "캠페인 비용이 커지면 precision이 높은 LR 또는 높은 threshold가 유리하다.",
        ],
        7.3,
        1.35,
        5.1,
        2.4,
        size=14,
    )
    add_footer(slide, 11)

    slide = blank_slide(prs)
    add_title(slide, "Calibration: 점수는 확률이 아니다", "Raw score를 이탈 확률로 말하면 안 된다")
    add_picture_fit(slide, PHASE6 / "phase6_calibration_comparison.png", 0.85, 1.2, 6.4, 4.4)
    add_table(
        slide,
        [
            ["Case", "Raw mean", "Platt mean", "Observed"],
            ["LR", "0.346", "0.065", "0.0646"],
            ["BalancedBagging", "0.467", "0.065", "0.0646"],
            ["CatBoost", "0.414", "0.065", "0.0646"],
            ["XGBoost", "0.469", "0.065", "0.0646"],
        ],
        7.65,
        1.45,
        4.8,
        1.75,
    )
    add_bullets(
        slide,
        [
            "Platt calibration 후 평균 score가 실제 churn rate와 맞춰졌다.",
            "Calibration은 probability 해석을 개선하지만 ranking 자체를 크게 바꾸지는 않는다.",
        ],
        7.75,
        3.75,
        4.6,
        1.2,
        size=13,
    )
    add_footer(slide, 12)

    slide = blank_slide(prs)
    add_title(slide, "Segment별 ROI와 실패 패턴", "전체 평균은 고객군별 운영 리스크를 숨긴다")
    add_picture_fit(slide, PHASE6 / "phase6_segment_recall_heatmap.png", 0.75, 1.15, 6.5, 4.7)
    add_bullets(
        slide,
        [
            "low-value에서는 CatBoost 계열이 많은 이탈자를 포착했다.",
            "mid-value에서는 BalancedBagging과 XGBoost가 강했다.",
            "high-value는 precision이 낮아 고객 피로도 관리가 중요하다.",
        ],
        7.7,
        1.45,
        4.8,
        2.2,
        size=14,
    )
    add_footer(slide, 13)

    slide = blank_slide(prs)
    add_title(slide, "모델 합의도 기반 Risk Tiering", "여러 모델이 동시에 경고한 고객은 더 위험하다")
    add_picture_fit(slide, PHASE6 / "phase6_model_agreement.png", 0.8, 1.2, 6.3, 4.6)
    add_bullets(
        slide,
        [
            "8개 모델 모두가 경고한 고객군의 이탈률은 12.43%.",
            "전체 평균 6.46%의 약 1.9배다.",
            "합의도는 campaign priority tier로 활용 가능하다.",
            "단, vote 0에서도 실제 이탈자가 있어 완전한 rule은 아니다.",
        ],
        7.65,
        1.35,
        4.8,
        2.8,
        size=14,
    )
    add_footer(slide, 14)

    slide = blank_slide(prs)
    add_title(slide, "통계 검정 보강", "Hold-out point estimate의 불확실성을 같이 제시")
    metric_ci = pd.read_csv(PHASE8 / "bootstrap_metric_ci.csv")
    rows = [["Case", "F1", "F1 CI", "Recall", "Recall CI"]]
    for case in ["LR_no_zip_f1", "BalancedBagging_with_zip", "CatBoost_native_with_zip", "XGBoost_with_zip"]:
        row = metric_ci[metric_ci["case_label"].eq(case)].iloc[0]
        rows.append(
            [
                case.replace("_with_zip", "").replace("_no_zip_f1", " LR"),
                f"{row['f1_point']:.3f}",
                f"[{row['f1_ci95_low']:.3f}, {row['f1_ci95_high']:.3f}]",
                f"{row['recall_point']:.3f}",
                f"[{row['recall_ci95_low']:.3f}, {row['recall_ci95_high']:.3f}]",
            ]
        )
    add_table(slide, rows, 0.7, 1.25, 12.0, 2.1)
    add_bullets(
        slide,
        [
            "LR F1 point는 가장 높지만 CI가 넓다.",
            "BalancedBagging/XGBoost의 recall 장점은 bootstrap CI에서도 유지된다.",
            "McNemar test는 주요 모델 쌍의 error pattern이 서로 다름을 보여준다.",
        ],
        1.0,
        4.25,
        11.5,
        1.3,
        size=14,
    )
    add_footer(slide, 15)

    slide = blank_slide(prs)
    add_title(slide, "한계와 향후 개선", "성능 상한의 핵심은 모델보다 데이터 구조")
    add_table(
        slide,
        [
            ["한계", "향후 개선"],
            ["정적 CRM snapshot", "월별 사용량/매출 추세 추가"],
            ["시간 기반 행동 feature 부재", "결제 실패, 문의, 계약 만료 정보 추가"],
            ["배포용 확률 보정 미반영", "calibrated probability 기반 threshold 재설계"],
            ["외부 검증 데이터 없음", "기간별 hold-out 또는 다른 시장 데이터 검증"],
        ],
        0.85,
        1.2,
        11.6,
        2.2,
    )
    add_callout(slide, "성능 개선의 다음 열쇠는 모델 추가가 아니라 temporal feature 확보", 1.1, 4.65, 11.0, 0.72, ACCENT_3)
    add_footer(slide, 16)

    slide = blank_slide(prs)
    add_band(slide, 0, 0, 13.333, 7.5, RGBColor(236, 242, 248))
    add_text(slide, "최종 결론", 0.8, 0.9, 11.5, 0.65, size=34, bold=True, color=TITLE)
    add_bullets(
        slide,
        [
            "논문 EasyEnsemble baseline은 F1 0.128로 재현했다.",
            "Hold-out F1 최고는 LogisticRegression_SMOTE지만, CV 안정성은 BalancedBagging/EasyEnsemble이 강하다.",
            "Recall 극대화, 비용 최적화, 예산 제한에 따라 최적 운영점은 달라진다.",
            "ChurnRadar의 핵심 성과는 단일 모델 승리가 아니라 목적별 운영 프레임워크다.",
        ],
        1.0,
        2.1,
        10.8,
        3.2,
        size=17,
    )
    add_footer(slide, 17)

    return prs


def main() -> None:
    prs = build_presentation()
    prs.save(OUTPUT)
    print(f"Saved: {OUTPUT}")


if __name__ == "__main__":
    main()
